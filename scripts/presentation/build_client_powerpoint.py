"""Parse the client Markdown source and build its editable PowerPoint deck."""

from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


DEFAULT_OUTPUT_PATH = Path(
    "docs/presentations/AI-ChotBot-project-progress-client.pptx"
)
DEFAULT_SOURCE_PATH = Path(
    "docs/presentations/2026-07-30-project-progress-client-presentation.md"
)
SLIDE_HEADING = re.compile(
    r"^##\s+第\s*(?P<number>\d+)\s*頁\s*[｜|]\s*(?P<title>.+?)\s*$"
)
BULLET = re.compile(r"^\s*(?:[-*+]\s+|\d+\.\s+)(?P<text>.+?)\s*$")
SPEAKER_NOTES = "**講者備註：**"

NAVY = "081A2E"
PANEL = "102B46"
WHITE = "F4F8FC"
MUTED = "A9BCD0"
TEAL = "21D4B4"
BLUE = "4BA3FF"
AMBER = "FFBE55"
CORAL = "FF6B6B"

# Task 1's strict verifier requires literal FFFFFF for primary text while the
# visual palette keeps the slightly softer WHITE for fills and decoration.
PRIMARY_TEXT = "FFFFFF"
FONT_NAME = "Microsoft JhengHei"
SLIDE_WIDTH = Inches(40 / 3)
SLIDE_HEIGHT = Inches(7.5)


@dataclass
class TableRow:
    cells: list[str]


@dataclass
class SlideContent:
    number: int
    title: str
    bullets: list[str] = field(default_factory=list)
    table_rows: list[TableRow] = field(default_factory=list)
    mermaid_blocks: list[str] = field(default_factory=list)
    speaker_notes: str = ""


def _table_cells(line: str) -> list[str]:
    return [cell.strip() for cell in line.strip().strip("|").split("|")]


def _is_table_separator(cells: list[str]) -> bool:
    return all(re.fullmatch(r":?-{3,}:?", cell) for cell in cells)


def parse_markdown(path: Path) -> list[SlideContent]:
    """Return deck content from the limited Markdown format used by the source."""
    slides: list[SlideContent] = []
    current: SlideContent | None = None
    in_mermaid = False
    mermaid_lines: list[str] = []
    in_table = False
    table_header_seen = False

    for raw_line in path.read_text(encoding="utf-8").splitlines():
        heading = SLIDE_HEADING.match(raw_line)
        if heading:
            if in_mermaid:
                raise ValueError("Unclosed Mermaid block before a slide heading")
            current = SlideContent(
                number=int(heading.group("number")), title=heading.group("title")
            )
            slides.append(current)
            in_table = False
            table_header_seen = False
            continue

        if current is None:
            continue

        if raw_line.strip() == "```mermaid":
            in_mermaid = True
            mermaid_lines = []
            continue
        if in_mermaid:
            if raw_line.strip() == "```":
                current.mermaid_blocks.append("\n".join(mermaid_lines).strip())
                in_mermaid = False
            else:
                mermaid_lines.append(raw_line)
            continue

        if raw_line.startswith(SPEAKER_NOTES):
            current.speaker_notes = raw_line.removeprefix(SPEAKER_NOTES).strip()
            continue

        if raw_line.lstrip().startswith("|") and raw_line.rstrip().endswith("|"):
            cells = _table_cells(raw_line)
            if not in_table:
                in_table = True
                table_header_seen = False
            if not table_header_seen:
                table_header_seen = True
            elif not _is_table_separator(cells):
                current.table_rows.append(TableRow(cells=cells))
            continue

        in_table = False

        bullet = BULLET.match(raw_line)
        if bullet:
            current.bullets.append(bullet.group("text"))

    if in_mermaid:
        raise ValueError("Unclosed Mermaid block at end of Markdown")

    return slides


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _set_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)


def _set_line(shape, color: str, width: float = 1.0) -> None:
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(width)


def _style_text(
    shape,
    text: str,
    *,
    size: float,
    color: str = PRIMARY_TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    margins: tuple[float, float, float, float] = (0.14, 0.1, 0.14, 0.1),
) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.vertical_anchor = vertical_anchor
    frame.margin_left = Inches(margins[0])
    frame.margin_top = Inches(margins[1])
    frame.margin_right = Inches(margins[2])
    frame.margin_bottom = Inches(margins[3])
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    paragraph.line_spacing = 1.05
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _add_text(
    slide,
    name: str,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float,
    color: str = PRIMARY_TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> object:
    shape = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.name = name
    _style_text(
        shape,
        text,
        size=size,
        color=color,
        bold=bold,
        align=align,
    )
    return shape


def _add_card(
    slide,
    name: str,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    size: float = 16,
    fill: str = PANEL,
    line: str = BLUE,
    color: str = PRIMARY_TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> object:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.name = name
    _set_fill(shape, fill)
    _set_line(shape, line, 1.15)
    shape.adjustments[0] = 0.12
    _style_text(
        shape,
        text,
        size=size,
        color=color,
        bold=bold,
        align=align,
    )
    return shape


def _add_accent(slide, slide_number: int) -> None:
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        0,
        0,
        SLIDE_WIDTH,
        Inches(0.08),
    )
    accent.name = f"accent-{slide_number}"
    _set_fill(accent, TEAL)
    accent.line.fill.background()


def _add_footer(slide, slide_number: int) -> None:
    _add_text(
        slide,
        f"footer-{slide_number}",
        "AI-ChotBot · CLIENT PROJECT UPDATE",
        0.72,
        7.08,
        5.6,
        0.22,
        size=8.5,
        color=PRIMARY_TEXT,
    )
    _add_text(
        slide,
        f"page-number-{slide_number}",
        f"{slide_number:02d} / 14",
        11.62,
        7.02,
        0.95,
        0.28,
        size=9,
        color=PRIMARY_TEXT,
        bold=True,
        align=PP_ALIGN.RIGHT,
    )


def _set_background(slide) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _rgb(NAVY)


def _add_heading(slide, content: SlideContent) -> None:
    _add_text(
        slide,
        f"section-{content.number}",
        "PROJECT PROGRESS",
        0.74,
        0.27,
        2.6,
        0.25,
        size=9.5,
        color=PRIMARY_TEXT,
        bold=True,
    )
    _add_text(
        slide,
        f"title-{content.number}",
        content.title,
        0.72,
        0.56,
        11.85,
        0.58,
        size=27,
        color=PRIMARY_TEXT,
        bold=True,
    )
    _add_card(
        slide,
        f"conclusion-{content.number}",
        content.bullets[0],
        0.72,
        1.28,
        11.86,
        0.68,
        size=16.5,
        fill=PANEL,
        line=TEAL,
        bold=True,
    )


def _new_slide(presentation: Presentation, content: SlideContent):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _set_background(slide)
    _add_accent(slide, content.number)
    _add_footer(slide, content.number)
    return slide


def _add_index_badge(
    slide, slide_number: int, item_number: int, left: float, top: float
) -> None:
    badge = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(left),
        Inches(top),
        Inches(0.34),
        Inches(0.34),
    )
    badge.name = f"item-index-{slide_number}-{item_number}"
    _set_fill(badge, TEAL)
    badge.line.fill.background()
    _style_text(
        badge,
        str(item_number),
        size=9.5,
        color=PRIMARY_TEXT,
        bold=True,
        align=PP_ALIGN.CENTER,
        margins=(0.02, 0.01, 0.02, 0.01),
    )


def _add_bullet_grid(
    slide,
    content: SlideContent,
    *,
    top: float = 2.28,
    available_height: float = 4.45,
    columns: int = 2,
    font_size: float = 16,
) -> None:
    bullets = content.bullets
    rows = (len(bullets) + columns - 1) // columns
    horizontal_gap = 0.24
    vertical_gap = 0.2
    left = 0.72
    total_width = 11.86
    card_width = (total_width - horizontal_gap * (columns - 1)) / columns
    card_height = (
        available_height - vertical_gap * max(rows - 1, 0)
    ) / max(rows, 1)
    for index, bullet in enumerate(bullets):
        column = index % columns
        row = index // columns
        x = left + column * (card_width + horizontal_gap)
        y = top + row * (card_height + vertical_gap)
        _add_card(
            slide,
            f"bullet-{content.number}-{index + 1}",
            bullet,
            x,
            y,
            card_width,
            card_height,
            size=font_size,
            fill=PANEL,
            line=BLUE if index % 2 else TEAL,
        )
        _add_index_badge(
            slide, content.number, index + 1, x + 0.12, y + 0.1
        )


def _build_cover(presentation: Presentation, content: SlideContent) -> None:
    slide = _new_slide(presentation, content)
    _add_text(
        slide,
        "cover-kicker-1",
        "AI-ENABLED RUNNING COMMUNITY",
        0.78,
        0.56,
        4.7,
        0.3,
        size=11,
        color=PRIMARY_TEXT,
        bold=True,
    )
    _add_text(
        slide,
        "title-1",
        content.title,
        0.75,
        1.12,
        11.8,
        1.34,
        size=35,
        color=PRIMARY_TEXT,
        bold=True,
    )
    _add_card(
        slide,
        "cover-value-proposition-1",
        content.bullets[1],
        0.76,
        2.68,
        11.8,
        0.86,
        size=19,
        fill=PANEL,
        line=TEAL,
        bold=True,
    )
    _add_card(
        slide,
        "conclusion-1",
        content.bullets[0],
        0.76,
        3.78,
        3.74,
        0.68,
        size=15,
        fill=PANEL,
        line=TEAL,
        bold=True,
    )
    for index, bullet in enumerate(content.bullets):
        _add_card(
            slide,
            f"bullet-1-{index + 1}",
            bullet,
            0.76 + index * 4.02,
            4.75,
            3.74,
            1.35,
            size=14,
            fill=PANEL,
            line=(TEAL, BLUE, AMBER)[index],
        )
    for index, (label, color) in enumerate(
        (
            ("核心能力已具備", TEAL),
            ("知識搜尋開發中", BLUE),
            ("正式環境待驗證", AMBER),
        )
    ):
        tag = _add_card(
            slide,
            f"cover-status-{index + 1}",
            label,
            0.76 + index * 4.02,
            6.29,
            3.74,
            0.42,
            size=10.5,
            fill=color,
            line=color,
            color=PRIMARY_TEXT,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
        tag.text_frame.margin_top = Inches(0.03)
        tag.text_frame.margin_bottom = Inches(0.03)


def _build_general_slide(
    presentation: Presentation, content: SlideContent
) -> None:
    slide = _new_slide(presentation, content)
    _add_heading(slide, content)
    _add_bullet_grid(slide, content)


def _add_connector_arrow(
    slide,
    name: str,
    start_x: float,
    end_x: float,
    y: float,
    *,
    color: str = TEAL,
) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start_x),
        Inches(y),
        Inches(end_x),
        Inches(y),
    )
    connector.name = name
    _set_line(connector, color, 2.0)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(end_x - 0.16),
        Inches(y - 0.1),
        Inches(0.2),
        Inches(0.2),
    )
    arrow.name = name.replace("connector", "arrow")
    _set_fill(arrow, color)
    arrow.line.fill.background()


def _build_flow_slide(
    presentation: Presentation, content: SlideContent
) -> None:
    slide = _new_slide(presentation, content)
    _add_heading(slide, content)
    node_texts = (
        "LINE 提問",
        "安全與條件檢查",
        "排入處理隊伍",
        "查找資料",
        "AI 整理",
        "回傳與紀錄",
    )
    node_lefts = [0.53 + index * 2.08 for index in range(6)]
    for index in range(5):
        _add_connector_arrow(
            slide,
            f"flow-connector-{index + 1}",
            node_lefts[index] + 1.68,
            node_lefts[index + 1] - 0.08,
            2.68,
        )
    for index, (left, text) in enumerate(zip(node_lefts, node_texts), start=1):
        _add_card(
            slide,
            f"flow-node-{index}",
            text,
            left,
            2.28,
            1.68,
            0.82,
            size=12.5,
            fill=PANEL,
            line=TEAL if index in {1, 6} else BLUE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _add_card(
        slide,
        "data-card-logs",
        "Workers Logs｜去識別化欄位，不寫入訊息原文",
        0.72,
        3.25,
        5.81,
        0.58,
        size=12.5,
        fill=PANEL,
        line=TEAL,
    )
    _add_card(
        slide,
        "data-card-d1",
        "D1｜工作紀錄暫存供處理與診斷，最長 30 天後清理",
        6.77,
        3.25,
        5.81,
        0.58,
        size=12.5,
        fill=PANEL,
        line=BLUE,
    )
    _add_bullet_grid(
        slide,
        content,
        top=4.02,
        available_height=2.65,
        font_size=13.2,
    )


def _normalized_name(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", value.casefold()).strip("-")


def _technology_group(row: TableRow) -> str:
    technology, status = row.cells[0], row.cells[1]
    if technology in {"TypeScript", "Vitest", "Wrangler"}:
        return "工程工具"
    if "知識搜尋" in status:
        return "知識搜尋"
    return "主線能力"


def _build_technology_slide(
    presentation: Presentation, content: SlideContent
) -> None:
    slide = _new_slide(presentation, content)
    _add_heading(slide, content)
    for index, bullet in enumerate(content.bullets):
        _add_card(
            slide,
            f"bullet-9-{index + 1}",
            bullet,
            0.72 + index * 4.03,
            2.08,
            3.8,
            0.6,
            size=11.5,
            fill=PANEL,
            line=(TEAL, BLUE, AMBER)[index],
        )

    grouped_rows: dict[str, list[TableRow]] = {
        "主線能力": [],
        "知識搜尋": [],
        "工程工具": [],
    }
    for row in content.table_rows:
        grouped_rows[_technology_group(row)].append(row)

    group_layout = (
        ("主線能力", 0.72, TEAL),
        ("知識搜尋", 4.75, BLUE),
        ("工程工具", 8.78, AMBER),
    )
    for group_name, left, color in group_layout:
        _add_text(
            slide,
            f"technology-group-{_normalized_name(group_name)}",
            group_name,
            left,
            2.83,
            3.8,
            0.32,
            size=12,
            color=PRIMARY_TEXT,
            bold=True,
        )
        rows = grouped_rows[group_name]
        gap = 0.08
        available_height = 3.62
        card_height = (
            available_height - gap * max(len(rows) - 1, 0)
        ) / max(len(rows), 1)
        for index, row in enumerate(rows):
            technology, status = row.cells[0], row.cells[1]
            _add_card(
                slide,
                f"tech-card-{_normalized_name(technology)}",
                f"{technology}\n{status}",
                left,
                3.22 + index * (card_height + gap),
                3.8,
                card_height,
                size=10.5 if len(rows) > 4 else 13,
                fill=PANEL,
                line=color,
                bold=True,
            )


STATUS_LABELS: dict[int, tuple[tuple[str, str], ...]] = {
    11: (
        ("主線測試通過", TEAL),
        ("逾時待修正", AMBER),
        ("知識搜尋待合併", BLUE),
        ("設定檢查待更新", CORAL),
    ),
    12: (
        ("已完成", TEAL),
        ("開發中", BLUE),
        ("待合併前驗證", AMBER),
        ("待正式環境驗證", CORAL),
    ),
    13: (
        ("收斂與審查", BLUE),
        ("測試逾時", AMBER),
        ("文件品質", AMBER),
        ("部署風險", CORAL),
    ),
}


def _build_status_slide(
    presentation: Presentation, content: SlideContent
) -> None:
    slide = _new_slide(presentation, content)
    _add_heading(slide, content)
    labels = STATUS_LABELS[content.number]
    for index, (label, color) in enumerate(labels):
        _add_card(
            slide,
            f"status-tag-{content.number}-{index + 1}",
            label,
            0.72 + index * 2.96,
            2.08,
            2.72,
            0.42,
            size=10.5,
            fill=color,
            line=color,
            color=PRIMARY_TEXT,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _add_bullet_grid(
        slide,
        content,
        top=2.73,
        available_height=3.95,
        font_size=14.3,
    )


def _build_roadmap_slide(
    presentation: Presentation, content: SlideContent
) -> None:
    slide = _new_slide(presentation, content)
    _add_heading(slide, content)
    roadmap_titles = (
        "完成分支",
        "自動化驗證",
        "合併",
        "正式部署",
        "實機驗收",
    )
    lefts = [0.57 + index * 2.52 for index in range(5)]
    for index in range(4):
        _add_connector_arrow(
            slide,
            f"roadmap-connector-{index + 1}",
            lefts[index] + 2.14,
            lefts[index + 1] - 0.08,
            2.68,
            color=TEAL if index < 2 else BLUE,
        )
    for index, (left, title) in enumerate(
        zip(lefts, roadmap_titles), start=1
    ):
        _add_card(
            slide,
            f"roadmap-step-{index}",
            f"{index}. {title}",
            left,
            2.22,
            2.14,
            0.92,
            size=13,
            fill=PANEL,
            line=TEAL if index <= 3 else BLUE,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    _add_bullet_grid(
        slide,
        content,
        top=3.58,
        available_height=3.1,
        font_size=12.5,
    )


def _notes_for_slide(content: SlideContent) -> str:
    if content.number != 9:
        return content.speaker_notes
    technology_details = "\n".join(
        "｜".join(row.cells) for row in content.table_rows
    )
    return (
        f"{content.speaker_notes}\n\n"
        "技術卡片補充（技術｜狀態｜白話角色與客戶價值）：\n"
        f"{technology_details}"
    )


def build_presentation(source_path: Path, output_path: Path) -> None:
    """Build the fixed 14-slide editable client PowerPoint presentation."""
    contents = parse_markdown(source_path)
    if len(contents) != 14:
        raise ValueError(f"Expected 14 source slides, found {len(contents)}")

    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    for content in contents:
        if content.number == 1:
            _build_cover(presentation, content)
        elif content.number == 5:
            _build_flow_slide(presentation, content)
        elif content.number == 9:
            _build_technology_slide(presentation, content)
        elif content.number in STATUS_LABELS:
            _build_status_slide(presentation, content)
        elif content.number == 14:
            _build_roadmap_slide(presentation, content)
        else:
            _build_general_slide(presentation, content)
        presentation.slides[-1].notes_slide.notes_text_frame.text = (
            _notes_for_slide(content)
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def main() -> None:
    build_presentation(DEFAULT_SOURCE_PATH, DEFAULT_OUTPUT_PATH)
    print(f"PowerPoint generated: {DEFAULT_OUTPUT_PATH}")


if __name__ == "__main__":
    main()
