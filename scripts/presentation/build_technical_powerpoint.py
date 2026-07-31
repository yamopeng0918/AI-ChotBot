"""Build the editable 16-slide technical achievements PowerPoint."""

from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

try:
    from .build_client_powerpoint import parse_markdown
except ImportError:  # Supports direct execution from scripts/presentation.
    from build_client_powerpoint import parse_markdown


DEFAULT_SOURCE_PATH = Path(
    "docs/presentations/2026-07-30-technical-achievements-presentation.md"
)
DEFAULT_OUTPUT_PATH = Path(
    "docs/presentations/AI-ChotBot-technical-achievements.pptx"
)

NAVY = "081A2E"
PANEL = "102B46"
WHITE = "F4F8FC"
MUTED = "A9BCD0"
TEAL = "21D4B4"
BLUE = "4BA3FF"
AMBER = "FFBE55"
CORAL = "FF6B6B"
FONT_NAME = "Microsoft JhengHei"


def _display_text(text: str) -> str:
    """Render the limited inline Markdown used in source bullets as plain text."""
    return re.sub(r"(\*\*|__|`)", "", text).strip()


def _rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def _set_fill(shape, color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)


def _set_line(shape, color: str, width: float = 1.0) -> None:
    shape.line.color.rgb = _rgb(color)
    shape.line.width = Pt(width)


def _style(shape, text: str, *, size: float, color: str = WHITE, bold: bool = False, align=PP_ALIGN.LEFT) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.margin_top = Inches(0.07)
    frame.margin_bottom = Inches(0.07)
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = _rgb(color)


def _text(slide, name: str, text: str, left: float, top: float, width: float, height: float, *, size: float = 16, color: str = WHITE, bold: bool = False, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    shape.name = name
    _style(shape, text, size=size, color=color, bold=bold, align=align)
    return shape


def _card(slide, name: str, text: str, left: float, top: float, width: float, height: float, *, line: str = BLUE, fill: str = PANEL, size: float = 15, bold: bool = False, color: str = WHITE, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.name = name
    _set_fill(shape, fill)
    _set_line(shape, line, 1.1)
    shape.adjustments[0] = 0.1
    _style(shape, text, size=size, color=color, bold=bold, align=align)
    return shape


def _connector(slide, name: str, start_x: float, start_y: float, end_x: float, end_y: float, *, color: str = TEAL):
    shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(start_x), Inches(start_y), Inches(end_x), Inches(end_y),
    )
    shape.name = name
    _set_line(shape, color, 1.8)
    arrow = OxmlElement("a:tailEnd")
    arrow.set("type", "triangle")
    shape._element.spPr.ln.append(arrow)
    return shape


def _base_slide(presentation: Presentation, source) -> object:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = _rgb(NAVY)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(40 / 3), Inches(0.08))
    accent.name = f"accent-{source.number}"
    _set_fill(accent, TEAL)
    accent.line.fill.background()
    _text(slide, f"section-{source.number}", "AI-CHOTBOT · TECHNICAL ACHIEVEMENTS", 0.7, 0.24, 5.0, 0.22, size=8.5, color=TEAL, bold=True)
    _text(slide, f"title-{source.number}", source.title, 0.7, 0.52, 12.0, 0.5, size=25, bold=True)
    _card(slide, f"conclusion-{source.number}", _display_text(source.bullets[0]), 0.7, 1.16, 11.95, 0.58, line=TEAL, size=16, bold=True)
    _text(slide, f"footer-{source.number}", "16 頁技術成果與工程治理簡報", 0.7, 7.08, 5.0, 0.18, size=8, color=MUTED)
    _text(slide, f"page-number-{source.number}", f"{source.number:02d} / 16", 11.7, 7.04, 0.8, 0.2, size=8.5, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    slide.notes_slide.notes_text_frame.text = source.speaker_notes
    return slide


def _source_bullets(slide, source, *, top: float = 2.05, compact: bool = False) -> None:
    columns = 2
    rows = (len(source.bullets) + columns - 1) // columns
    width = 5.82
    height = 0.66 if compact else (4.55 - 0.16 * (rows - 1)) / rows
    for index, bullet in enumerate(source.bullets, start=1):
        column = (index - 1) % columns
        row = (index - 1) // columns
        _card(
            slide,
            f"bullet-{source.number}-{index}",
            _display_text(bullet),
            0.7 + column * 6.12,
            top + row * (height + (0.08 if compact else 0.16)),
            width,
            height,
            line=BLUE if index % 2 == 0 else TEAL,
            size=16,
        )


def _architecture(slide) -> None:
    labels = ("LINE", "Webhook", "Worker", "Queue", "Workers AI", "D1", "Open-Meteo")
    for index, label in enumerate(labels, start=1):
        column = (index - 1) % 4
        row = (index - 1) // 4
        _card(slide, f"architecture-node-{index}", f"架構節點 {index}\n{label}", 0.75 + column * 3.1, 2.15 + row * 1.05, 2.45, 0.64, line=TEAL if row == 0 else BLUE, size=13, bold=True, align=PP_ALIGN.CENTER)
    points = ((3.2, 2.47, 3.85, 2.47), (6.3, 2.47, 6.95, 2.47), (9.4, 2.47, 10.05, 2.47), (4.0, 2.82, 2.2, 3.2), (7.0, 2.82, 5.3, 3.2), (10.0, 2.82, 8.4, 3.2))
    for index, point in enumerate(points, start=1):
        _connector(slide, f"architecture-connector-{index}", *point)


def _message_flow(slide) -> None:
    labels = ("Webhook 接收", "簽章與群組 mention 驗證", "Queue 入列", "AI 或資料來源處理", "LINE reply 與 push fallback", "D1 與觀測紀錄")
    for index, label in enumerate(labels):
        column = index % 3
        row = index // 3
        _card(slide, f"message-flow-step-{index + 1}", label, 0.85 + column * 4.05, 2.2 + row * 1.0, 2.5, 0.52, line=BLUE, size=12.5, bold=True, align=PP_ALIGN.CENTER)
    points = ((3.35, 2.46, 4.9, 2.46), (7.4, 2.46, 8.95, 2.46), (10.2, 2.72, 10.2, 3.2), (8.95, 3.46, 7.4, 3.46), (4.9, 3.46, 3.35, 3.46))
    for index, point in enumerate(points, start=1):
        _connector(slide, f"message-flow-connector-{index}", *point, color=BLUE)


def _worker_flow(slide) -> None:
    labels = ("Webhook 路由", "Hono", "Queue consumer", "綁定服務")
    for index, label in enumerate(labels, start=1):
        _card(slide, f"worker-flow-node-{index}", label, 0.9 + (index - 1) * 3.05, 2.4, 2.25, 0.58, line=TEAL, size=13.5, bold=True, align=PP_ALIGN.CENTER)
        if index > 1:
            _connector(slide, f"worker-flow-connector-{index - 1}", 0.9 + (index - 2) * 3.05 + 2.25, 2.69, 0.9 + (index - 1) * 3.05, 2.69)


def _reliability(slide) -> None:
    nodes = (("retry", "retry\n暫時失敗再處理", TEAL), ("dlq", "DLQ\n反覆失敗待調查", CORAL), ("deduplication", "deduplication\n避免重複回覆", BLUE))
    for index, (name, text, color) in enumerate(nodes):
        _card(slide, f"reliability-node-{name}", text, 1.0 + index * 4.05, 2.3, 3.0, 0.8, line=color, size=14, bold=True, align=PP_ALIGN.CENTER)
        if index:
            _connector(slide, f"reliability-connector-{index}", 1.0 + (index - 1) * 4.05 + 3.0, 2.7, 1.0 + index * 4.05, 2.7, color=color)


def _model_boundary(slide) -> None:
    values = (("primary", "主要模型", TEAL), ("fallback", "備援模型", BLUE), ("policy", "回答政策", AMBER))
    for index, (name, text, color) in enumerate(values):
        _card(slide, f"model-boundary-{name}", text, 0.9 + index * 4.1, 2.35, 3.05, 0.68, line=color, size=14, bold=True, align=PP_ALIGN.CENTER)


def _weather_cache_flow(slide) -> None:
    values = (("intent", "意圖辨識"), ("cache-read", "快取讀取"), ("provider", "Open-Meteo"), ("cache-write", "快取寫入"))
    for index, (name, text) in enumerate(values):
        _card(slide, f"weather-cache-step-{name}", text, 0.85 + index * 3.1, 2.35, 2.45, 0.62, line=TEAL if index < 2 else BLUE, size=13.5, bold=True, align=PP_ALIGN.CENTER)
        if index:
            _connector(slide, f"weather-cache-connector-{index}", 0.85 + (index - 1) * 3.1 + 2.45, 2.66, 0.85 + index * 3.1, 2.66, color=BLUE)


def _data_lifecycle(slide) -> None:
    values = (("question-record", "D1 問題紀錄"), ("weather-cache", "weather cache"), ("group-settings", "group settings"), ("metrics", "metrics"), ("lifecycle", "30-day lifecycle"))
    for index, (name, text) in enumerate(values):
        _card(slide, f"data-node-{name}", text, 0.8 + (index % 3) * 4.05, 2.05 + (index // 3) * 0.85, 3.25, 0.53, line=BLUE if index % 2 else TEAL, size=13, bold=True, align=PP_ALIGN.CENTER)


def _privacy_layers(slide) -> None:
    values = (("logs", "Workers Logs", TEAL), ("d1", "D1", BLUE), ("secrets", "Cloudflare secrets", AMBER))
    for index, (name, text, color) in enumerate(values):
        _card(slide, f"privacy-layer-{name}", text, 0.9 + index * 4.1, 2.35, 3.05, 0.68, line=color, size=14, bold=True, align=PP_ALIGN.CENTER)


def _observability(slide) -> None:
    _card(slide, "observability-correlation-webhook", "webhookEventId", 0.85, 2.05, 2.4, 0.5, line=TEAL, size=13, bold=True, align=PP_ALIGN.CENTER)
    _card(slide, "observability-correlation-operation", "operationId", 3.45, 2.05, 2.4, 0.5, line=BLUE, size=13, bold=True, align=PP_ALIGN.CENTER)
    events = ("webhook.enqueue.completed", "question.started", "storage.claim.completed", "answer.completed", "line.reply.completed")
    for index, event in enumerate(events):
        _card(slide, f"observability-event-{index + 1}", event, 0.85 + index * 2.38, 3.05, 2.0, 0.5, line=AMBER, size=12, align=PP_ALIGN.CENTER)
        if index:
            _connector(slide, f"observability-connector-{index}", 0.85 + (index - 1) * 2.38 + 2.0, 3.3, 0.85 + index * 2.38, 3.3, color=AMBER)


def _quality_gates(slide) -> None:
    values = (("main", "主線 134 通過", TEAL), ("knowledge", "知識搜尋 421 通過", BLUE), ("timeout", "1 項超過 5 秒", CORAL), ("predeploy", "設定檢查待更新", AMBER))
    for index, (name, text, color) in enumerate(values):
        _card(slide, f"quality-gate-{name}", text, 0.9 + (index % 2) * 6.0, 2.15 + (index // 2) * 0.9, 5.55, 0.6, line=color, size=14, bold=True, align=PP_ALIGN.CENTER)


def _knowledge_architecture(slide) -> None:
    names = ("r2", "ingestion-queue", "workers-ai", "vectorize", "retrieval", "grounded-answer")
    for index, name in enumerate(names):
        _card(slide, f"knowledge-node-{name}", name, 0.85 + (index % 3) * 4.05, 2.0 + (index // 3) * 0.88, 3.1, 0.54, line=BLUE, size=13, bold=True, align=PP_ALIGN.CENTER)
    _card(slide, "development-status-13", "開發中", 10.5, 3.1, 1.6, 0.46, line=AMBER, fill=AMBER, size=13, bold=True, align=PP_ALIGN.CENTER)


def _maturity_matrix(slide) -> None:
    _card(slide, "maturity-matrix-14", "已完成｜開發中｜待合併前驗證｜待正式環境驗證", 0.85, 2.12, 11.55, 0.68, line=TEAL, size=15, bold=True, align=PP_ALIGN.CENTER)


def _roadmap(slide, source) -> None:
    for index, bullet in enumerate(source.bullets, start=1):
        number, body = bullet.split(". ", 1)
        label = f"{number}. {body.split('：', 1)[0]}"
        _card(slide, f"roadmap-step-{index}", label, 0.85 + (index - 1) * 2.45, 2.1, 2.15, 1.0, line=TEAL if index < 4 else AMBER, size=16, bold=True, align=PP_ALIGN.CENTER)
        if index > 1:
            _connector(slide, f"roadmap-connector-{index - 1}", 0.85 + (index - 2) * 2.45 + 2.15, 2.6, 0.85 + (index - 1) * 2.45, 2.6, color=TEAL if index < 4 else AMBER)


def build_technical_presentation(source_path: Path, output_path: Path) -> None:
    """Generate an editable technical deck from the 16-slide Markdown source."""
    source_slides = parse_markdown(source_path)
    if [slide.number for slide in source_slides] != list(range(1, 17)):
        raise ValueError("Technical source must contain slides 1 through 16")

    presentation = Presentation()
    presentation.slide_width = Inches(40 / 3)
    presentation.slide_height = Inches(7.5)

    special_numbers = {3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 16}
    for source in source_slides:
        slide = _base_slide(presentation, source)
        _source_bullets(slide, source, top=4.65 if source.number in special_numbers else 2.05, compact=source.number in special_numbers)
        if source.number == 3:
            _architecture(slide)
        elif source.number == 4:
            _message_flow(slide)
        elif source.number == 5:
            _worker_flow(slide)
        elif source.number == 6:
            _reliability(slide)
        elif source.number == 7:
            _model_boundary(slide)
        elif source.number == 8:
            _weather_cache_flow(slide)
        elif source.number == 9:
            _data_lifecycle(slide)
        elif source.number == 10:
            _privacy_layers(slide)
        elif source.number == 11:
            _observability(slide)
        elif source.number == 12:
            _quality_gates(slide)
        elif source.number == 13:
            _knowledge_architecture(slide)
        elif source.number == 14:
            _maturity_matrix(slide)
        elif source.number == 16:
            _roadmap(slide, source)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def main() -> None:
    build_technical_presentation(DEFAULT_SOURCE_PATH, DEFAULT_OUTPUT_PATH)
    print(f"Technical PowerPoint generated: {DEFAULT_OUTPUT_PATH}")


if __name__ == "__main__":
    main()
