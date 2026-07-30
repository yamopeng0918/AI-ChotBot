"""Validate a generated client PowerPoint against its Markdown source."""

from __future__ import annotations

from pathlib import Path
import re

from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

try:
    from .build_client_powerpoint import STATUS_LABELS, parse_markdown, roadmap_label
except ImportError:  # Direct execution: python scripts/presentation/verify_*.py
    from build_client_powerpoint import STATUS_LABELS, parse_markdown, roadmap_label


SENSITIVE_IDENTIFIER = re.compile(
    r"access[_ -]?token|channel[_ -]?secret|analytics_hash_key|database_id",
    re.IGNORECASE,
)
EXPECTED_BACKGROUND = "081A2E"
EXPECTED_PRIMARY_TEXT = "F4F8FC"
EXPECTED_ACCENT = "21D4B4"
EXPECTED_FONT = "Microsoft JhengHei"
DECORATIVE_TEXT_PREFIXES = (
    "date-",
    "footer-",
    "page-number-",
    "slide-number-",
)


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _shape_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in _iter_shapes(slide.shapes):
        if getattr(shape, "has_table", False):
            texts.extend(
                cell.text.strip()
                for row in shape.table.rows
                for cell in row.cells
                if cell.text.strip()
            )
        elif getattr(shape, "has_text_frame", False) and shape.text.strip():
            texts.append(shape.text.strip())
    return texts


def _named_shapes(slide, prefix: str) -> list:
    normalized_prefix = prefix.casefold()
    return [
        shape
        for shape in _iter_shapes(slide.shapes)
        if shape.name.casefold().startswith(normalized_prefix)
    ]


def _is_native_text_shape(shape) -> bool:
    return (
        shape.shape_type != MSO_SHAPE_TYPE.PICTURE
        and getattr(shape, "has_text_frame", False)
        and bool(shape.text.strip())
    )


def _normalized_name(value: str) -> str:
    return re.sub(r"[^a-z0-9]+", "-", value.casefold()).strip("-")


def _notes_text(slide) -> str:
    return slide.notes_slide.notes_text_frame.text.strip()


def _shape_is_outside_slide(shape, slide_width: int, slide_height: int) -> bool:
    return (
        shape.left < 0
        or shape.top < 0
        or shape.left + shape.width > slide_width
        or shape.top + shape.height > slide_height
    )


def _rgb_value(color_format) -> str | None:
    try:
        if color_format.type == MSO_COLOR_TYPE.RGB:
            return str(color_format.rgb).upper()
    except (AttributeError, TypeError, ValueError):
        pass
    return None


def _slide_background_rgb(slide) -> str | None:
    fill = slide.background.fill
    if fill.type != MSO_FILL_TYPE.SOLID:
        return None
    return _rgb_value(fill.fore_color)


def _shape_runs(shape):
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                for paragraph in cell.text_frame.paragraphs:
                    yield from paragraph.runs
    elif getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            yield from paragraph.runs


def _is_decorative_text_shape(shape) -> bool:
    if shape.name.casefold().startswith(DECORATIVE_TEXT_PREFIXES):
        return True
    if not getattr(shape, "is_placeholder", False):
        return False
    return shape.placeholder_format.type in {
        PP_PLACEHOLDER.DATE,
        PP_PLACEHOLDER.FOOTER,
        PP_PLACEHOLDER.SLIDE_NUMBER,
    }


def _shape_is_visible(shape, slide_width: int, slide_height: int) -> bool:
    return (
        shape.width > 0
        and shape.height > 0
        and shape.left < slide_width
        and shape.top < slide_height
        and shape.left + shape.width > 0
        and shape.top + shape.height > 0
    )


def _primary_text_style_counts(
    slide, slide_width: int, slide_height: int
) -> tuple[int, int, int]:
    total_weight = 0
    white_weight = 0
    expected_font_weight = 0
    for shape in _iter_shapes(slide.shapes):
        if (
            _is_decorative_text_shape(shape)
            or not _shape_is_visible(shape, slide_width, slide_height)
        ):
            continue
        for run in _shape_runs(shape):
            weight = len(run.text.strip())
            if not weight:
                continue
            total_weight += weight
            if _rgb_value(run.font.color) == EXPECTED_PRIMARY_TEXT:
                white_weight += weight
            if (run.font.name or "").casefold() == EXPECTED_FONT.casefold():
                expected_font_weight += weight
    return total_weight, white_weight, expected_font_weight


def _has_minimum_declared_font_size(shape, minimum_points: float) -> bool:
    text_runs = [
        run
        for run in _shape_runs(shape)
        if run.text.strip()
    ]
    return bool(text_runs) and all(
        run.font.size is not None and run.font.size.pt >= minimum_points
        for run in text_runs
    )


def _has_expected_accent(
    slide, slide_width: int, slide_height: int
) -> bool:
    for shape in _iter_shapes(slide.shapes):
        if (
            shape.width <= 0
            or shape.height <= 0
            or shape.left >= slide_width
            or shape.top >= slide_height
            or shape.left + shape.width <= 0
            or shape.top + shape.height <= 0
        ):
            continue
        try:
            if (
                shape.fill.type == MSO_FILL_TYPE.SOLID
                and _rgb_value(shape.fill.fore_color) == EXPECTED_ACCENT
            ):
                return True
        except (AttributeError, TypeError, ValueError):
            pass
        try:
            if _rgb_value(shape.line.color) == EXPECTED_ACCENT:
                return True
        except (AttributeError, TypeError, ValueError):
            pass
        for run in _shape_runs(shape):
            if _rgb_value(run.font.color) == EXPECTED_ACCENT:
                return True
    return False


def _validate_editable_named_shapes(
    slide,
    prefix: str,
    expected_count: int,
    slide_number: int,
) -> tuple[list, list[str]]:
    shapes = _named_shapes(slide, prefix)
    errors: list[str] = []
    if len(shapes) != expected_count:
        errors.append(
            f"Slide {slide_number} must contain exactly {expected_count} "
            f"{prefix} shapes"
        )
    if any(not _is_native_text_shape(shape) for shape in shapes):
        errors.append(
            f"Slide {slide_number} {prefix} items must be native text shapes"
        )
    return shapes, errors


def verify_presentation(pptx_path: Path, source_path: Path) -> list[str]:
    """Return contract violations for ``pptx_path``; an empty list means success."""
    presentation = Presentation(pptx_path)
    source_slides = parse_markdown(source_path)
    errors: list[str] = []

    source_numbers = [slide.number for slide in source_slides]
    if len(source_slides) != 14:
        errors.append(
            f"Source must contain exactly 14 slides, found {len(source_slides)}"
        )
    if source_numbers != list(range(1, 15)):
        errors.append("Source slide numbers must be exactly 1..14 in order")
    if len(presentation.slides) != 14:
        errors.append(f"Expected 14 slides, found {len(presentation.slides)}")
    if abs(presentation.slide_width * 9 - presentation.slide_height * 16) > 10:
        errors.append("Presentation aspect ratio must be 16:9")

    for index, source_slide in enumerate(source_slides):
        if index >= len(presentation.slides):
            break
        slide = presentation.slides[index]
        slide_number = index + 1
        texts = _shape_texts(slide)
        notes = _notes_text(slide)

        total_text, white_text, expected_font_text = _primary_text_style_counts(
            slide, presentation.slide_width, presentation.slide_height
        )
        if not total_text or white_text / total_text < 0.75:
            errors.append(
                f"Slide {slide_number}: at least 75% of visible primary text "
                f"must be white ({EXPECTED_PRIMARY_TEXT})"
            )
        if not total_text or expected_font_text / total_text < 0.75:
            errors.append(
                f"Slide {slide_number}: at least 75% of visible primary text "
                f"must use {EXPECTED_FONT}"
            )
        if _slide_background_rgb(slide) != EXPECTED_BACKGROUND:
            errors.append(
                f"Slide {slide_number} background must be {EXPECTED_BACKGROUND}"
            )
        if not _has_expected_accent(
            slide, presentation.slide_width, presentation.slide_height
        ):
            errors.append(
                f"Slide {slide_number} must include a visible "
                f"{EXPECTED_ACCENT} teal accent"
            )
        if source_slide.title not in texts:
            errors.append(f"Slide {slide_number} title does not match the source")
        if not notes:
            errors.append(f"Slide {slide_number} speaker notes are empty")
        elif source_slide.speaker_notes not in notes:
            errors.append(
                f"Slide {slide_number} notes must contain the source speaker notes"
            )

        conclusions = _named_shapes(slide, "conclusion-")
        if len(conclusions) != 1 or not _is_native_text_shape(conclusions[0]):
            errors.append(
                f"Slide {slide_number} must contain exactly one native "
                "conclusion- text shape"
            )
        elif conclusions[0].text.strip() != source_slide.bullets[0]:
            errors.append(
                f"Slide {slide_number} conclusion must equal the source conclusion"
            )
        bullets = _named_shapes(slide, "bullet-")
        if not 3 <= len(bullets) <= 5:
            errors.append(
                f"Slide {slide_number} must contain 3–5 bullet- shapes"
            )
        if any(not _is_native_text_shape(shape) for shape in bullets):
            errors.append(
                f"Slide {slide_number} bullet- items must be native text shapes"
            )
        elif [shape.text.strip() for shape in bullets] != source_slide.bullets:
            errors.append(
                f"Slide {slide_number} bullet text and order must match "
                "the source bullets"
            )
        for prefix, minimum_points in (
            ("conclusion-", 16),
            ("bullet-", 16),
            ("cover-value-proposition-", 16),
            ("tech-card-", 12),
            ("roadmap-step-", 16),
        ):
            for shape in _named_shapes(slide, prefix):
                if not _has_minimum_declared_font_size(shape, minimum_points):
                    errors.append(
                        f"Slide {slide_number} {shape.name} text must declare "
                        f"a font size of at least {minimum_points}pt"
                    )

        for shape in _iter_shapes(slide.shapes):
            if _shape_is_outside_slide(
                shape, presentation.slide_width, presentation.slide_height
            ):
                errors.append(f"Slide {slide_number} has a shape outside slide bounds")

        sensitive_match = SENSITIVE_IDENTIFIER.search("\n".join([*texts, notes]))
        if sensitive_match:
            errors.append(
                f"Sensitive identifier found on slide {slide_number}: "
                f"{sensitive_match.group(0)}"
            )

        if slide_number == 5:
            flow_nodes = _named_shapes(slide, "flow-node-")
            if len(flow_nodes) < 6:
                errors.append(
                    "Slide 5 must contain at least 6 flow-node- shapes"
                )
            if any(not _is_native_text_shape(shape) for shape in flow_nodes):
                errors.append(
                    "Slide 5 flow-node- items must be native text shapes"
                )
            expected_flow_labels = re.findall(
                r"\b[A-F]\[([^\]]+)\]",
                source_slide.mermaid_blocks[0],
            )
            if [
                shape.text.strip()
                for shape in flow_nodes
                if _is_native_text_shape(shape)
            ] != expected_flow_labels:
                errors.append(
                    "Slide 5 flow-node text must match every source flow node"
                )
        if slide_number == 9:
            tech_cards, tech_errors = _validate_editable_named_shapes(
                slide, "tech-card-", 13, slide_number
            )
            errors.extend(tech_errors)
            for row in source_slide.table_rows:
                technology, status, value = row.cells[:3]
                normalized_shape_name = f"tech-card-{_normalized_name(technology)}"
                matches = [
                    shape
                    for shape in tech_cards
                    if _is_native_text_shape(shape)
                    and shape.name.casefold() == normalized_shape_name
                ]
                if len(matches) != 1 or technology.casefold() not in (
                    matches[0].text.casefold() if matches else ""
                ):
                    errors.append(
                        "Slide 9 must contain every source technical name "
                        f"exactly once; missing or duplicated: {technology}"
                    )
                    continue
                if status.casefold() not in matches[0].text.casefold():
                    errors.append(
                        f"Slide 9 technical status must match the source for "
                        f"{technology}"
                    )
                visible_or_notes = "\n".join([*texts, notes]).casefold()
                if value.casefold() not in visible_or_notes:
                    errors.append(
                        f"Slide 9 technical value must be visible or in notes "
                        f"for {technology}"
                    )
        if slide_number in STATUS_LABELS:
            expected_statuses = STATUS_LABELS[slide_number]
            status_shapes, status_errors = _validate_editable_named_shapes(
                slide, "status-tag-", len(expected_statuses), slide_number
            )
            errors.extend(status_errors)
            for index, (expected_label, _) in enumerate(
                expected_statuses, start=1
            ):
                expected_name = f"status-tag-{slide_number}-{index}"
                matches = [
                    shape
                    for shape in status_shapes
                    if shape.name.casefold() == expected_name
                ]
                if (
                    len(matches) != 1
                    or matches[0].text.strip() != expected_label
                ):
                    errors.append(
                        f"Slide {slide_number} {expected_name} must match "
                        "the configured visible status label"
                    )
        if slide_number == 14:
            roadmap_shapes = _named_shapes(slide, "roadmap-step-")
            expected_names = [
                f"roadmap-step-{number}" for number in range(1, 6)
            ]
            for expected_name in expected_names:
                matches = [
                    shape
                    for shape in roadmap_shapes
                    if shape.name.casefold() == expected_name
                ]
                if len(matches) != 1:
                    errors.append(
                        f"Slide 14 must contain {expected_name} exactly once"
                    )
                elif not _is_native_text_shape(matches[0]):
                    errors.append(
                        f"Slide 14 {expected_name} must be a native text shape"
                    )
                elif matches[0].text.strip() != roadmap_label(
                    source_slide.bullets[
                        int(expected_name.rsplit("-", 1)[1]) - 1
                    ]
                ):
                    errors.append(
                        f"Slide 14 {expected_name} must match the compact "
                        "source roadmap label"
                    )
            if len(roadmap_shapes) != 5:
                errors.append(
                    "Slide 14 must contain exactly five roadmap-step- shapes"
                )

    return errors


def main() -> None:
    pptx_path = Path(
        "docs/presentations/AI-ChotBot-project-progress-client.pptx"
    )
    source_path = Path(
        "docs/presentations/2026-07-30-project-progress-client-presentation.md"
    )
    errors = verify_presentation(pptx_path, source_path)
    if errors:
        for error in errors:
            print(f"ERROR: {error}")
        raise SystemExit(1)
    print("PowerPoint verification passed: 14 slides")


if __name__ == "__main__":
    main()
