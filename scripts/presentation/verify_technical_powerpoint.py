"""Verify the 16-slide technical deck against its editable source contract."""

from __future__ import annotations

from pathlib import Path
import re
import sys
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Inches

try:
    from .build_client_powerpoint import parse_markdown
except ImportError:  # Supports direct execution from scripts/presentation.
    from build_client_powerpoint import parse_markdown


DEFAULT_SOURCE_PATH = Path("docs/presentations/2026-07-30-technical-achievements-presentation.md")
DEFAULT_PPTX_PATH = Path("docs/presentations/AI-ChotBot-technical-achievements.pptx")
HAN = re.compile(r"[\u4e00-\u9fff]")
SENSITIVE = re.compile(r"access[_ -]?token|channel[_ -]?secret|analytics_hash_key|database_id", re.I)
NAVY = "081A2E"
WHITE = "F4F8FC"
TEAL = "21D4B4"
FONT_NAME = "Microsoft JhengHei"
FOOTER_SAFE_TOP = Inches(7.0)
DECORATIVE_PREFIXES = ("accent-", "footer-", "page-number-", "section-")
TECH_CARD_PREFIXES = (
    "architecture-node-", "message-flow-step-", "worker-flow-node-",
    "reliability-node-", "model-boundary-", "weather-cache-step-",
    "data-node-", "privacy-layer-", "observability-event-",
    "quality-gate-", "knowledge-node-", "development-status-",
    "maturity-matrix-",
)


def _iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes(shape.shapes)


def _validate_ooxml(pptx_path: Path) -> list[str]:
    """Check ZIP integrity and reject packaged content outside this deck contract."""
    errors: list[str] = []
    try:
        with zipfile.ZipFile(pptx_path) as archive:
            bad_member = archive.testzip()
            if bad_member:
                errors.append(f"OOXML CRC check failed for {bad_member}")
            names = archive.namelist()
            lowered = [name.casefold() for name in names]
            if any(name.startswith("ppt/media/") for name in lowered):
                errors.append("OOXML package must not contain image or media parts")
            if any("vbaproject" in name for name in lowered):
                errors.append("OOXML package must not contain VBA parts")
            for name in names:
                if not name.casefold().endswith((".xml", ".rels")):
                    continue
                try:
                    root = etree.fromstring(archive.read(name))
                except (etree.XMLSyntaxError, zipfile.BadZipFile) as error:
                    errors.append(f"OOXML XML cannot be parsed: {name}: {error}")
                    continue
                if name.casefold().endswith(".rels"):
                    for relationship in root:
                        if relationship.get("TargetMode") == "External":
                            errors.append(f"OOXML external relationship is forbidden: {name}")
    except (OSError, zipfile.BadZipFile) as error:
        errors.append(f"OOXML ZIP cannot be opened: {error}")
    return errors


def _shape_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in _iter_shapes(slide.shapes):
        if getattr(shape, "has_table", False):
            texts.extend(cell.text.strip() for row in shape.table.rows for cell in row.cells if cell.text.strip())
        elif getattr(shape, "has_text_frame", False) and shape.text.strip():
            texts.append(shape.text.strip())
    return texts


def _notes_text(slide) -> str:
    return slide.notes_slide.notes_text_frame.text.strip()


def _native_text(shape) -> bool:
    return shape.shape_type != MSO_SHAPE_TYPE.PICTURE and getattr(shape, "has_text_frame", False) and bool(shape.text.strip())


def _rgb_value(color_format) -> str | None:
    try:
        return str(color_format.rgb).upper()
    except (AttributeError, TypeError, ValueError):
        return None


def _has_primary_style(shape) -> bool:
    if not _native_text(shape):
        return False
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue
            if run.font.name != FONT_NAME or run.font.size is None or run.font.size.pt < 16:
                return False
            if _rgb_value(run.font.color) != WHITE:
                return False
    return True


def _has_minimum_font_size(shape, minimum_points: float) -> bool:
    if not _native_text(shape):
        return False
    runs = [run for paragraph in shape.text_frame.paragraphs for run in paragraph.runs if run.text]
    return bool(runs) and all(run.font.size is not None and run.font.size.pt >= minimum_points for run in runs)


def _display_text(text: str) -> str:
    return re.sub(r"(\*\*|__|`)", "", text).strip()


def _text_density(shape) -> float:
    area = shape.width.inches * shape.height.inches
    return len(shape.text.replace("\n", "")) / area if area else float("inf")


def _require_readable_card(shape, name: str, errors: list[str], *, minimum_height: float) -> None:
    if shape.height.inches < minimum_height:
        errors.append(f"{name} needs at least {minimum_height:.2f} inches of internal height")
    if _text_density(shape) > 25:
        errors.append(f"{name} exceeds the 25 chars/in² density limit")


def _has_triangle_tail(shape) -> bool:
    if shape.shape_type != MSO_SHAPE_TYPE.LINE:
        return False
    line = shape._element.spPr.ln
    tail = line.find(qn("a:tailEnd"))
    return tail is not None and tail.get("type") == "triangle"


def _named(slide, prefix: str) -> list:
    return [shape for shape in _iter_shapes(slide.shapes) if shape.name.startswith(prefix)]


def _require_exact_text(slide, name: str, expected: str, errors: list[str], slide_number: int, label: str) -> None:
    expected = _display_text(expected)
    shapes = _named(slide, name)
    exact = [shape for shape in shapes if shape.name == name]
    if len(shapes) != 1 or len(exact) != 1:
        errors.append(f"Slide {slide_number} requires {name} exactly once (no copy suffix)")
    elif not _native_text(exact[0]):
        errors.append(f"Slide {slide_number} {name} must be a native non-empty text shape")
    elif exact[0].text.strip() != expected:
        errors.append(f"Slide {slide_number} {label} must match the source")
    elif name.startswith(("title-", "conclusion-", "bullet-", "roadmap-step-")) and not _has_primary_style(exact[0]):
        errors.append(f"Slide {slide_number} {name} must use {FONT_NAME}, {WHITE}, and at least 16pt")
    elif name.startswith(("conclusion-", "bullet-", "roadmap-step-")):
        _require_readable_card(exact[0], name, errors, minimum_height=0.55)


def _require_named_terms(slide, prefix: str, terms: tuple[str, ...], errors: list[str], slide_number: int, expected_count: int | None = None) -> None:
    shapes = _named(slide, prefix)
    expected_count = len(terms) if expected_count is None else expected_count
    if len(shapes) != expected_count or any(not _native_text(shape) for shape in shapes):
        errors.append(f"Slide {slide_number} {prefix} shapes must be native non-empty text shapes exactly once")
        return
    text = "\n".join(shape.text.strip().casefold() for shape in shapes)
    for term in terms:
        if term.casefold() not in text:
            errors.append(f"Slide {slide_number} {prefix} must include {term}")


def _require_named_keyword(slide, name: str, keyword: str, errors: list[str], slide_number: int) -> None:
    shapes = _named(slide, name)
    exact = [shape for shape in shapes if shape.name == name]
    if len(shapes) != 1 or len(exact) != 1:
        errors.append(f"Slide {slide_number} requires {name} exactly once (no copy suffix)")
    elif not _native_text(exact[0]):
        errors.append(f"Slide {slide_number} {name} must be a native non-empty text shape")
    elif keyword.casefold() not in exact[0].text.casefold():
        errors.append(f"Slide {slide_number} {name} must include {keyword}")


def _require_named_line(slide, name: str, errors: list[str]) -> None:
    shapes = _named(slide, name)
    exact = [shape for shape in shapes if shape.name == name]
    if len(shapes) != 1 or len(exact) != 1 or exact[0].shape_type != MSO_SHAPE_TYPE.LINE:
        errors.append(f"Slide 4 requires {name} exactly once as a native LINE connector")
    elif not _has_triangle_tail(exact[0]):
        errors.append(f"Slide 4 requires {name} to end with a triangle arrow")


def _require_arrow_connectors(slide, prefix: str, expected_count: int, errors: list[str], slide_number: int) -> None:
    shapes = _named(slide, prefix)
    if len(shapes) != expected_count or any(not _has_triangle_tail(shape) for shape in shapes):
        errors.append(f"Slide {slide_number} requires {expected_count} {prefix} connectors with triangle arrows")


def _require_left_to_right_connector(slide, name: str, errors: list[str], slide_number: int) -> None:
    shapes = _named(slide, name)
    exact = [shape for shape in shapes if shape.name == name]
    if len(shapes) != 1 or len(exact) != 1:
        errors.append(f"Slide {slide_number} requires {name} exactly once")
    elif not _has_triangle_tail(exact[0]) or exact[0].begin_x >= exact[0].end_x:
        errors.append(f"Slide {slide_number} {name} must be a left-to-right triangle-arrow connector")


def _roadmap_label(bullet: str) -> str:
    number, body = bullet.split(". ", 1)
    return f"{number}. {body.split('：', 1)[0]}"


def _validate_source(source_path: Path) -> tuple[list, list[str]]:
    try:
        slides = parse_markdown(source_path)
    except (OSError, ValueError) as error:
        return [], [f"Technical source cannot be parsed: {error}"]
    errors: list[str] = []
    if len(slides) != 16:
        return slides, [f"Technical source must contain 16 slides, found {len(slides)}"]
    if [slide.number for slide in slides] != list(range(1, 17)):
        errors.append("Technical source slide numbers must be 1..16 in order")
    for slide in slides:
        if not slide.title.strip() or not HAN.search(slide.title):
            errors.append(f"Technical source slide {slide.number} needs a Traditional Chinese title")
        if not 3 <= len(slide.bullets) <= 5:
            errors.append(f"Technical source slide {slide.number} must have 3–5 bullet points")
        if any(not bullet.strip() or not HAN.search(bullet) for bullet in slide.bullets):
            errors.append(f"Technical source slide {slide.number} needs Traditional Chinese bullets")
        if not slide.speaker_notes.strip() or not HAN.search(slide.speaker_notes):
            errors.append(f"Technical source slide {slide.number} needs Traditional Chinese speaker notes")
    if len(slides) == 16:
        if "開發中" not in " ".join(slides[12].bullets):
            errors.append("Technical source slide 13 must explicitly state 開發中")
        if "待合併前驗證" not in " ".join(slides[13].bullets):
            errors.append("Technical source slide 14 must state 待合併前驗證")
    return slides, errors


def verify_technical_presentation(pptx_path: Path, source_path: Path) -> list[str]:
    """Return every source, editable-text, and special-page contract violation."""
    source_slides, errors = _validate_source(source_path)
    if not pptx_path.is_file():
        return [*errors, f"Technical PowerPoint does not exist: {pptx_path}"]
    ooxml_errors = _validate_ooxml(pptx_path)
    if any(error.startswith(("OOXML ZIP", "OOXML CRC", "OOXML XML")) for error in ooxml_errors):
        return [*errors, *ooxml_errors]
    errors.extend(ooxml_errors)
    try:
        presentation = Presentation(pptx_path)
    except Exception as error:  # pragma: no cover
        return [*errors, f"Technical PowerPoint cannot be opened: {error}"]
    if len(presentation.slides) != 16:
        return [*errors, f"Technical PowerPoint must contain 16 slides, found {len(presentation.slides)}"]

    if abs(presentation.slide_width.inches - (40 / 3)) > 0.001 or presentation.slide_height != Inches(7.5):
        errors.append("Technical PowerPoint must use exact 16:9 dimensions (13.333 × 7.5 inches)")
    for index, slide in enumerate(presentation.slides, start=1):
        try:
            background = _rgb_value(slide.background.fill.fore_color)
        except AttributeError:
            background = None
        if background != NAVY:
            errors.append(f"Slide {index} must use deep navy {NAVY} background")
        accents = _named(slide, f"accent-{index}")
        exact_accents = [shape for shape in accents if shape.name == f"accent-{index}"]
        if len(accents) != 1 or len(exact_accents) != 1 or _rgb_value(exact_accents[0].fill.fore_color) != TEAL:
            errors.append(f"Slide {index} needs one visible teal {TEAL} accent")
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                errors.append(f"Slide {index} must not contain images")
            if shape.name.startswith(DECORATIVE_PREFIXES):
                continue
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > presentation.slide_width or shape.top + shape.height > presentation.slide_height:
                errors.append(f"Slide {index} shape {shape.name} is outside the slide bounds")
            elif shape.top + shape.height > FOOTER_SAFE_TOP:
                errors.append(f"Slide {index} shape {shape.name} overlaps the footer-safe area")
            if shape.name.startswith(TECH_CARD_PREFIXES):
                if not _has_minimum_font_size(shape, 12):
                    errors.append(f"Slide {index} {shape.name} must declare at least 12pt text")
                else:
                    _require_readable_card(shape, shape.name, errors, minimum_height=0.45)

    for source in source_slides:
        slide = presentation.slides[source.number - 1]
        _require_exact_text(slide, f"title-{source.number}", source.title, errors, source.number, "title")
        _require_exact_text(slide, f"conclusion-{source.number}", source.bullets[0], errors, source.number, "conclusion")
        bullet_shapes = _named(slide, f"bullet-{source.number}-")
        if len(bullet_shapes) != len(source.bullets):
            errors.append(f"Slide {source.number} requires source bullet shapes exactly once")
        for index, bullet in enumerate(source.bullets, start=1):
            _require_exact_text(slide, f"bullet-{source.number}-{index}", bullet, errors, source.number, "bullet")
        notes = _notes_text(slide)
        if source.speaker_notes not in notes:
            errors.append(f"Slide {source.number} speaker notes must contain the source speaker notes")
        sensitive = SENSITIVE.search("\n".join([*_shape_texts(slide), notes]))
        if sensitive:
            errors.append(f"Sensitive identifier found on slide {source.number}: {sensitive.group(0)}")

    if len(source_slides) != 16:
        return errors
    slide = presentation.slides[2]
    architecture_roles = ("使用者", "LINE", "Worker", "Queue", "AI", "D1", "Open-Meteo")
    for index, role in enumerate(architecture_roles, start=1):
        name = f"architecture-node-{index}"
        _require_exact_text(slide, name, role, errors, 3, "architecture-node")
        shapes = _named(slide, name)
        exact = [shape for shape in shapes if shape.name == name]
        if len(exact) == 1 and not _has_minimum_font_size(exact[0], 16):
            errors.append(f"Slide 3 {name} must declare at least 16pt text")
    connectors = _named(slide, "architecture-connector-")
    if len(connectors) < 6 or any(shape.shape_type != MSO_SHAPE_TYPE.LINE for shape in connectors):
        errors.append("Slide 3 needs at least six native architecture connectors")
    else:
        _require_arrow_connectors(slide, "architecture-connector-", len(connectors), errors, 3)
    flow_labels = ("Webhook 接收", "簽章與群組 mention 驗證", "Queue 入列", "AI 或資料來源處理", "LINE reply 與 push fallback", "D1 與觀測紀錄")
    for index, label in enumerate(flow_labels, start=1):
        _require_exact_text(presentation.slides[3], f"message-flow-step-{index}", label, errors, 4, "message-flow-step")
    for connector in range(1, 6):
        _require_named_line(presentation.slides[3], f"message-flow-connector-{connector}", errors)
    _require_arrow_connectors(presentation.slides[4], "worker-flow-connector-", 3, errors, 5)
    for name, keyword in (("retry", "retry"), ("dlq", "dlq"), ("deduplication", "deduplication")):
        _require_named_keyword(presentation.slides[5], f"reliability-node-{name}", keyword, errors, 6)
    _require_arrow_connectors(presentation.slides[5], "reliability-connector-", 2, errors, 6)
    for name, label in (("primary", "主要模型"), ("fallback", "備援模型"), ("policy", "回答政策")):
        _require_exact_text(presentation.slides[6], f"model-boundary-{name}", label, errors, 7, "model-boundary")
    for name, label in (("intent", "意圖辨識"), ("cache-read", "快取讀取"), ("provider", "Open-Meteo"), ("cache-write", "快取寫入")):
        _require_exact_text(presentation.slides[7], f"weather-cache-step-{name}", label, errors, 8, "weather-cache-step")
    _require_arrow_connectors(presentation.slides[7], "weather-cache-connector-", 3, errors, 8)
    for name, keyword in (("question-record", "問題紀錄"), ("weather-cache", "weather cache"), ("group-settings", "group settings"), ("metrics", "metrics"), ("lifecycle", "lifecycle")):
        _require_named_keyword(presentation.slides[8], f"data-node-{name}", keyword, errors, 9)
    for name, label in (("logs", "Workers Logs"), ("d1", "D1"), ("secrets", "Cloudflare secrets")):
        _require_exact_text(presentation.slides[9], f"privacy-layer-{name}", label, errors, 10, "privacy-layer")
    events = ("webhook.enqueue.completed", "question.started", "storage.claim.completed", "answer.completed", "line.reply.completed")
    for index, event in enumerate(events, start=1):
        _require_exact_text(presentation.slides[10], f"observability-event-{index}", event, errors, 11, "observability-event")
    _require_arrow_connectors(presentation.slides[10], "observability-connector-", 4, errors, 11)
    visible_11 = "\n".join(_shape_texts(presentation.slides[10]))
    if "webhookEventId" not in visible_11 or "operationId" not in visible_11:
        errors.append("Slide 11 needs webhookEventId and operationId native text")
    quality = "\n".join(_shape_texts(presentation.slides[11]))
    for pattern, label in ((r"主線\s*134\s*通過", "主線 134 通過"), (r"知識搜尋\s*421\s*通過", "知識搜尋 421 通過"), (r"1 項.*超過\s*5 秒", "1 項超過 5 秒"), (r"設定檢查待更新", "設定檢查待更新")):
        if not re.search(pattern, quality):
            errors.append(f"Slide 12 needs quality-gate evidence: {label}")
    knowledge_names = ("r2", "ingestion-queue", "workers-ai", "vectorize", "retrieval", "grounded-answer")
    knowledge_lefts: list[int] = []
    for name in knowledge_names:
        _require_named_keyword(presentation.slides[12], f"knowledge-node-{name}", name, errors, 13)
        shapes = _named(presentation.slides[12], f"knowledge-node-{name}")
        exact = [shape for shape in shapes if shape.name == f"knowledge-node-{name}"]
        if len(exact) == 1:
            knowledge_lefts.append(exact[0].left)
    if knowledge_lefts != sorted(knowledge_lefts):
        errors.append("Slide 13 knowledge-node order must be R2 → ingestion queue → Workers AI → Vectorize → retrieval → grounded answer")
    for index in range(1, 6):
        _require_left_to_right_connector(presentation.slides[12], f"knowledge-connector-{index}", errors, 13)
    _require_exact_text(presentation.slides[12], "development-status-13", "開發中", errors, 13, "development-status")
    _require_named_keyword(presentation.slides[13], "maturity-matrix-14", "已完成", errors, 14)
    maturity_text = "\n".join(
        shape.text.strip()
        for shape in _named(presentation.slides[13], "maturity-matrix-14")
        if shape.name == "maturity-matrix-14" and _native_text(shape)
    )
    for state in ("開發中", "待合併前驗證", "待正式環境驗證"):
        if state not in maturity_text:
            errors.append(f"Slide 14 maturity-matrix-14 must include {state}")
    if re.search(r"\b\d+(?:\.\d+)?\s*%", "\n".join(_shape_texts(presentation.slides[13]))):
        errors.append("Slide 14 must not use percentage maturity")
    for step, bullet in enumerate(source_slides[15].bullets, start=1):
        _require_exact_text(presentation.slides[15], f"roadmap-step-{step}", _roadmap_label(bullet), errors, 16, f"roadmap-step-{step}")
    _require_arrow_connectors(presentation.slides[15], "roadmap-connector-", 4, errors, 16)
    return errors


def main() -> None:
    errors = verify_technical_presentation(DEFAULT_PPTX_PATH, DEFAULT_SOURCE_PATH)
    if errors:
        print("\n".join(errors))
        raise SystemExit(1)
    print("Technical PowerPoint verification passed: 16 slides")


if __name__ == "__main__":
    main()
