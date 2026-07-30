"""Contract tests for the client presentation parser and verifier."""

from __future__ import annotations

import base64
from io import BytesIO
from pathlib import Path
import re
import subprocess
import sys
from tempfile import TemporaryDirectory
import unittest

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

from scripts.presentation import build_client_powerpoint as deck_builder
from scripts.presentation.build_client_powerpoint import (
    DEFAULT_OUTPUT_PATH,
    parse_markdown,
)
from scripts.presentation.verify_client_powerpoint import verify_presentation


REPOSITORY_ROOT = Path(__file__).resolve().parents[2]
SOURCE_PATH = (
    REPOSITORY_ROOT
    / "docs"
    / "presentations"
    / "2026-07-30-project-progress-client-presentation.md"
)


def _roadmap_label(bullet: str) -> str:
    match = re.match(r"^(?P<number>\d+)\.\s*(?P<title>[^：:]+)", bullet)
    if not match:
        raise AssertionError(f"Unexpected roadmap bullet: {bullet}")
    return f"{match.group('number')} {match.group('title').strip()}"


class MarkdownParserTests(unittest.TestCase):
    def test_default_output_path_is_the_client_delivery_path(self) -> None:
        self.assertEqual(
            Path("docs/presentations/AI-ChotBot-project-progress-client.pptx"),
            DEFAULT_OUTPUT_PATH,
        )

    def test_parses_all_fourteen_numbered_slides_with_titles_and_notes(self) -> None:
        slides = parse_markdown(SOURCE_PATH)

        self.assertEqual(14, len(slides))
        self.assertEqual(list(range(1, 15)), [slide.number for slide in slides])
        self.assertTrue(all(slide.title.strip() for slide in slides))
        self.assertTrue(all(slide.speaker_notes.strip() for slide in slides))

    def test_readme_documents_powerpoint_runtime_and_rendering_limit(self) -> None:
        readme = (REPOSITORY_ROOT / "README.md").read_text(encoding="utf-8")

        self.assertIn("Python 3.14", readme)
        self.assertIn("python-pptx 1.0.2", readme)
        self.assertIn("PowerPoint/LibreOffice", readme)
        self.assertIn("換行", readme)
        self.assertIn("重疊", readme)
        self.assertIn("字級", readme)

    def test_parses_the_thirteen_technical_table_rows_on_slide_nine(self) -> None:
        slides = parse_markdown(SOURCE_PATH)

        self.assertEqual(13, len(slides[8].table_rows))
        self.assertTrue(all(row.cells for row in slides[8].table_rows))

    def test_preserves_bullets_and_mermaid_blocks_without_markdown_engine(self) -> None:
        slides = parse_markdown(SOURCE_PATH)

        self.assertGreaterEqual(len(slides[4].bullets), 4)
        self.assertEqual(1, len(slides[4].mermaid_blocks))
        self.assertIn("flowchart", slides[4].mermaid_blocks[0])


class GeneratedPresentationLayoutTests(unittest.TestCase):
    STATUS_COLORS = {"21D4B4", "4BA3FF", "FFBE55", "FF6B6B"}

    def _build_deck(self, destination: Path) -> Presentation:
        build_presentation = getattr(deck_builder, "build_presentation", None)
        self.assertIsNotNone(
            build_presentation,
            "build_presentation must generate the editable client deck",
        )
        build_presentation(SOURCE_PATH, destination)
        return Presentation(destination)

    @staticmethod
    def _shape_fill_rgb(shape) -> str | None:
        try:
            return str(shape.fill.fore_color.rgb).upper()
        except (AttributeError, TypeError, ValueError):
            return None

    def test_uses_exact_widescreen_dimensions(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")

        self.assertAlmostEqual(13.333, deck.slide_width.inches, places=3)
        self.assertEqual(7.5, deck.slide_height.inches)

    def test_generated_deck_satisfies_complete_verifier_contract(self) -> None:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            deck_builder.build_presentation(SOURCE_PATH, pptx_path)

            self.assertEqual([], verify_presentation(pptx_path, SOURCE_PATH))

    def test_primary_text_uses_exact_theme_white(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")

        title = next(
            shape
            for shape in deck.slides[1].shapes
            if shape.name == "title-2"
        )
        self.assertEqual(
            "F4F8FC",
            str(title.text_frame.paragraphs[0].runs[0].font.color.rgb),
        )

    def test_cover_names_the_value_proposition(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")
            cover = deck.slides[0]

        value_shapes = [
            shape
            for shape in cover.shapes
            if shape.name.startswith("cover-value-proposition-")
        ]
        self.assertEqual(1, len(value_shapes))
        self.assertEqual(
            parse_markdown(SOURCE_PATH)[0].bullets[1],
            value_shapes[0].text.strip(),
        )

    def test_slide_five_uses_native_flow_connectors(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")
            flow_slide = deck.slides[4]

        connectors = [
            shape
            for shape in flow_slide.shapes
            if shape.name.startswith("flow-connector-")
        ]
        self.assertGreaterEqual(5, len(connectors))
        self.assertTrue(
            all(shape.shape_type == MSO_SHAPE_TYPE.LINE for shape in connectors)
        )

    def test_slide_five_flow_nodes_match_mermaid_source_labels(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")
            flow_slide = deck.slides[4]

        mermaid = parse_markdown(SOURCE_PATH)[4].mermaid_blocks[0]
        expected_labels = re.findall(r"\b[A-F]\[([^\]]+)\]", mermaid)
        actual_labels = [
            shape.text.strip()
            for shape in flow_slide.shapes
            if shape.name.startswith("flow-node-")
        ]
        self.assertEqual(expected_labels, actual_labels)

    def test_slide_nine_contains_thirteen_native_technology_cards(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")
            technology_slide = deck.slides[8]

        cards = [
            shape
            for shape in technology_slide.shapes
            if shape.name.startswith("tech-card-")
        ]
        self.assertEqual(13, len(cards))
        self.assertTrue(
            all(
                card.shape_type != MSO_SHAPE_TYPE.PICTURE
                and card.has_text_frame
                for card in cards
            )
        )
        source_rows = parse_markdown(SOURCE_PATH)[8].table_rows
        notes = technology_slide.notes_slide.notes_text_frame.text
        cards_by_name = {card.name: card for card in cards}
        for row in source_rows:
            card = cards_by_name[
                "tech-card-"
                + re.sub(
                    r"[^a-z0-9]+", "-", row.cells[0].casefold()
                ).strip("-")
            ]
            self.assertIn(row.cells[0], card.text)
            self.assertIn(row.cells[1], card.text)
            self.assertIn(row.cells[2], notes)

    def test_slide_nine_technology_group_names_are_unique_ascii_slugs(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")
            technology_slide = deck.slides[8]

        names = [
            shape.name
            for shape in technology_slide.shapes
            if shape.name.startswith("technology-group-")
        ]
        self.assertEqual(3, len(names))
        self.assertEqual(3, len(set(names)))
        self.assertTrue(all(re.fullmatch(r"[a-z0-9-]+", name) for name in names))

    def test_readability_shapes_use_declared_minimum_font_sizes(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")

        minimums = {
            "conclusion-": 16,
            "bullet-": 16,
            "cover-value-proposition-": 16,
            "tech-card-": 12,
            "roadmap-step-": 16,
        }
        for slide in deck.slides:
            for shape in slide.shapes:
                for prefix, minimum in minimums.items():
                    if not shape.name.startswith(prefix):
                        continue
                    sizes = [
                        run.font.size.pt
                        for paragraph in shape.text_frame.paragraphs
                        for run in paragraph.runs
                        if run.text.strip() and run.font.size is not None
                    ]
                    self.assertTrue(sizes, shape.name)
                    self.assertGreaterEqual(min(sizes), minimum, shape.name)

    def test_slides_eleven_to_thirteen_have_named_status_colors(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")

        for slide_number in range(11, 14):
            status_shapes = [
                shape
                for shape in deck.slides[slide_number - 1].shapes
                if shape.name.startswith("status-tag-")
            ]
            self.assertTrue(status_shapes, f"slide {slide_number}")
            self.assertTrue(
                any(
                    self._shape_fill_rgb(shape) in self.STATUS_COLORS
                    for shape in status_shapes
                ),
                f"slide {slide_number}",
            )

    def test_slide_fourteen_has_five_numbered_roadmap_steps(self) -> None:
        with TemporaryDirectory() as directory:
            deck = self._build_deck(Path(directory) / "deck.pptx")
            roadmap_slide = deck.slides[13]

        roadmap_steps = [
            shape
            for shape in roadmap_slide.shapes
            if shape.name.startswith("roadmap-step-")
        ]
        self.assertEqual(
            [f"roadmap-step-{number}" for number in range(1, 6)],
            [shape.name for shape in roadmap_steps],
        )
        self.assertEqual(
            [
                _roadmap_label(bullet)
                for bullet in parse_markdown(SOURCE_PATH)[13].bullets
            ],
            [shape.text.strip() for shape in roadmap_steps],
        )


class PresentationVerifierCommandTests(unittest.TestCase):
    def test_verifier_runs_from_the_documented_direct_script_command(self) -> None:
        completed = subprocess.run(
            [
                sys.executable,
                "scripts/presentation/verify_client_powerpoint.py",
            ],
            cwd=REPOSITORY_ROOT,
            capture_output=True,
            check=False,
        )

        self.assertEqual(
            0,
            completed.returncode,
            completed.stderr.decode(errors="replace"),
        )
        self.assertIn(
            b"PowerPoint verification passed: 14 slides",
            completed.stdout,
        )


class PresentationVerifierTests(unittest.TestCase):
    DARK_BLUE = "081A2E"
    WHITE = "F4F8FC"
    TEAL = "21D4B4"
    FONT_NAME = "Microsoft JhengHei"
    ONE_PIXEL_PNG = base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8"
        "/x8AAusB9Wl2nWQAAAAASUVORK5CYII="
    )

    @staticmethod
    def _normalized_name(value: str) -> str:
        return re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")

    def _add_text(
        self,
        shapes,
        name: str,
        text: str,
        left: float,
        top: float,
        width: float = 2.0,
        height: float = 0.3,
        *,
        color: str | None = None,
        font_name: str | None = None,
        font_size: float = 16,
    ):
        shape = shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )
        shape.name = name
        shape.text_frame.text = text
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor.from_string(color or self.WHITE)
                run.font.name = font_name or self.FONT_NAME
                run.font.size = Pt(font_size)
        return shape

    def _write_presentation(
        self,
        destination: Path,
        *,
        aspect_ratio: str = "wide",
        background_color: str | None = None,
        primary_text_color: str | None = None,
        primary_font_name: str | None = None,
        accent_color: str | None = None,
        missing_accent_slide: int | None = None,
        invisible_accent_slide: int | None = None,
        title_override: str | None = None,
        include_notes: bool = True,
        bullet_text_override: str | None = None,
        conclusion_text_override: str | None = None,
        notes_text_override: str | None = None,
        missing_conclusion_slide: int | None = None,
        bullet_count_by_slide: dict[int, int] | None = None,
        flow_node_count: int = 6,
        flow_node_text_override: str | None = None,
        technology_count: int = 13,
        technology_name_override: str | None = None,
        technology_status_override: str | None = None,
        technology_value_override: str | None = None,
        status_label_override: str | None = None,
        roadmap_names: list[str] | None = None,
        roadmap_text_override: str | None = None,
        replace_flow_node_with_picture: bool = False,
        replace_tech_card_with_picture: bool = False,
        replace_roadmap_step_with_picture: bool = False,
        include_sensitive_text: bool = False,
        include_grouped_sensitive_text: bool = False,
        include_grouped_sensitive_table: bool = False,
        include_out_of_bounds_shape: bool = False,
        include_unnamed_wrong_style: bool = False,
        wrong_style_slide: int | None = None,
        undersized_shape_name: str | None = None,
        source_slide_limit: int | None = None,
    ) -> None:
        source_slides = parse_markdown(SOURCE_PATH)
        presentation = Presentation()
        if aspect_ratio == "wide":
            presentation.slide_width = 13_333_333
            presentation.slide_height = 7_500_000

        for source_slide in source_slides[:source_slide_limit]:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = RGBColor.from_string(
                background_color or self.DARK_BLUE
            )
            title = (
                title_override
                if source_slide.number == 1 and title_override
                else source_slide.title
            )
            self._add_text(
                slide.shapes,
                f"title-{source_slide.number}",
                title,
                0.4,
                0.2,
                12,
                0.5,
                color=primary_text_color,
                font_name=primary_font_name,
                font_size=(
                    15
                    if undersized_shape_name == f"title-{source_slide.number}"
                    else 16
                ),
            )
            if source_slide.number != missing_conclusion_slide:
                self._add_text(
                    slide.shapes,
                    f"conclusion-{source_slide.number}",
                    (
                        conclusion_text_override
                        if source_slide.number == 2 and conclusion_text_override
                        else source_slide.bullets[0]
                    ),
                    0.5,
                    0.8,
                    8,
                    color=primary_text_color,
                    font_name=primary_font_name,
                    font_size=(
                        15
                        if undersized_shape_name
                        == f"conclusion-{source_slide.number}"
                        else 16
                    ),
                )
            bullets = list(source_slide.bullets)
            requested_count = (bullet_count_by_slide or {}).get(
                source_slide.number, len(bullets)
            )
            bullets = bullets[:requested_count]
            if requested_count > len(bullets):
                bullets.extend(
                    f"Extra key point {number}"
                    for number in range(len(bullets) + 1, requested_count + 1)
                )
            if source_slide.number == 2 and bullet_text_override:
                bullets[0] = bullet_text_override
            for bullet_index, bullet_text in enumerate(bullets):
                self._add_text(
                    slide.shapes,
                    f"bullet-{source_slide.number}-{bullet_index + 1}",
                    bullet_text,
                    0.5,
                    1.2 + bullet_index * 0.35,
                    4,
                    color=primary_text_color,
                    font_name=primary_font_name,
                    font_size=(
                        15
                        if undersized_shape_name
                        == f"bullet-{source_slide.number}-{bullet_index + 1}"
                        else 16
                    ),
                )
            if include_notes or source_slide.number != 2:
                notes_text = (
                    notes_text_override
                    if source_slide.number == 2 and notes_text_override
                    else source_slide.speaker_notes
                )
                if source_slide.number == 9:
                    detail_lines = []
                    for row_index, row in enumerate(source_slide.table_rows):
                        cells = list(row.cells)
                        if row_index == 0 and technology_value_override:
                            cells[2] = technology_value_override
                        detail_lines.append("｜".join(cells))
                    notes_text += "\n\n" + "\n".join(detail_lines)
                slide.notes_slide.notes_text_frame.text = notes_text
            if source_slide.number != missing_accent_slide:
                accent = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(
                        20
                        if source_slide.number == invisible_accent_slide
                        else 12
                    ),
                    Inches(0.2),
                    Inches(0.2),
                    Inches(0.2),
                )
                accent.name = f"accent-{source_slide.number}"
                accent.fill.solid()
                accent.fill.fore_color.rgb = RGBColor.from_string(
                    accent_color or self.TEAL
                )
            if source_slide.number == 5:
                for node in range(flow_node_count):
                    name = f"flow-node-{node + 1}"
                    if replace_flow_node_with_picture and node == 0:
                        shape = slide.shapes.add_picture(
                            BytesIO(self.ONE_PIXEL_PNG),
                            Inches(0.5 + node),
                            Inches(3),
                            Inches(0.8),
                            Inches(0.4),
                        )
                        shape.name = name
                    else:
                        self._add_text(
                            slide.shapes,
                            name,
                            (
                                flow_node_text_override
                                if node == 0 and flow_node_text_override
                                else re.findall(
                                    r"\b[A-F]\[([^\]]+)\]",
                                    source_slide.mermaid_blocks[0],
                                )[node]
                            ),
                            0.5 + node,
                            3,
                            0.8,
                            0.4,
                            color=primary_text_color,
                            font_name=primary_font_name,
                        )
            if source_slide.number == 9:
                for row_index, row in enumerate(
                    source_slide.table_rows[:technology_count]
                ):
                    technology = (
                        technology_name_override
                        if row_index == 0 and technology_name_override
                        else row.cells[0]
                    )
                    name = f"tech-card-{self._normalized_name(technology)}"
                    if replace_tech_card_with_picture and row_index == 0:
                        shape = slide.shapes.add_picture(
                            BytesIO(self.ONE_PIXEL_PNG),
                            Inches(0.5 + row_index * 0.3),
                            Inches(3),
                            Inches(0.25),
                            Inches(0.25),
                        )
                        shape.name = name
                    else:
                        status = (
                            technology_status_override
                            if row_index == 0 and technology_status_override
                            else row.cells[1]
                        )
                        self._add_text(
                            slide.shapes,
                            name,
                            f"{technology}\n{status}",
                            0.5 + (row_index % 4) * 3,
                            3 + (row_index // 4) * 0.5,
                            2.7,
                            0.35,
                            color=primary_text_color,
                            font_name=primary_font_name,
                            font_size=(
                                11
                                if undersized_shape_name == name
                                else 12
                            ),
                        )
            if source_slide.number == 14:
                for step_index, name in enumerate(
                    roadmap_names
                    or [f"roadmap-step-{number}" for number in range(1, 6)]
                ):
                    if replace_roadmap_step_with_picture and step_index == 0:
                        shape = slide.shapes.add_picture(
                            BytesIO(self.ONE_PIXEL_PNG),
                            Inches(0.5),
                            Inches(3 + step_index * 0.5),
                            Inches(1),
                            Inches(0.3),
                        )
                        shape.name = name
                    else:
                        self._add_text(
                            slide.shapes,
                            name,
                            (
                                roadmap_text_override
                                if step_index == 0 and roadmap_text_override
                                else _roadmap_label(
                                    source_slide.bullets[
                                        min(
                                            step_index,
                                            len(source_slide.bullets) - 1,
                                        )
                                    ]
                                )
                            ),
                            0.5,
                            3 + step_index * 0.5,
                            5,
                            0.3,
                            color=primary_text_color,
                            font_name=primary_font_name,
                            font_size=(
                                15
                                if undersized_shape_name == name
                                else 16
                            ),
                        )
            if source_slide.number in deck_builder.STATUS_LABELS:
                for label_index, (label, _) in enumerate(
                    deck_builder.STATUS_LABELS[source_slide.number], start=1
                ):
                    self._add_text(
                        slide.shapes,
                        f"status-tag-{source_slide.number}-{label_index}",
                        (
                            status_label_override
                            if source_slide.number == 11
                            and label_index == 1
                            and status_label_override
                            else label
                        ),
                        8,
                        3 + label_index * 0.4,
                        3,
                        0.3,
                    )

        first_slide = presentation.slides[0]
        if include_sensitive_text:
            self._add_text(
                first_slide.shapes, "extra-sensitive", "access_token", 0.5, 5
            )
        if include_grouped_sensitive_text:
            outer_group = first_slide.shapes.add_group_shape()
            inner_group = outer_group.shapes.add_group_shape()
            self._add_text(
                inner_group.shapes,
                "nested-sensitive",
                "channel_secret",
                0.5,
                5,
            )
        if include_grouped_sensitive_table:
            group = first_slide.shapes.add_group_shape()
            table_shape = first_slide.shapes.add_table(
                1, 1, Inches(0.5), Inches(5), Inches(2), Inches(0.4)
            )
            table_shape.table.cell(0, 0).text = "database_id"
            group.shapes._spTree.insert_element_before(
                table_shape._element, "p:extLst"
            )
        if include_out_of_bounds_shape:
            self._add_text(
                first_slide.shapes, "overflow", "Overflow", 14, 1, 1, 0.3
            )
        if include_unnamed_wrong_style:
            self._add_text(
                presentation.slides[1].shapes,
                "Visible body copy",
                "Unstyled visible content " * 100,
                0.5,
                5,
                10,
                1,
                color="000000",
                font_name="Arial",
            )
        if wrong_style_slide is not None:
            for shape in presentation.slides[wrong_style_slide - 1].shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = RGBColor.from_string("000000")
                        run.font.name = "Arial"
        presentation.save(destination)

    def _errors(self, **overrides) -> str:
        with TemporaryDirectory() as directory:
            pptx_path = Path(directory) / "deck.pptx"
            self._write_presentation(pptx_path, **overrides)
            return "\n".join(
                verify_presentation(pptx_path, SOURCE_PATH)
            ).lower()

    def test_accepts_a_complete_widescreen_presentation(self) -> None:
        self.assertEqual("", self._errors())

    def test_rejects_thirteen_slide_source(self) -> None:
        with TemporaryDirectory() as directory:
            source_path = Path(directory) / "source.md"
            source_text = SOURCE_PATH.read_text(encoding="utf-8")
            source_path.write_text(
                source_text.rsplit("\n---\n", 1)[0],
                encoding="utf-8",
            )
            pptx_path = Path(directory) / "deck.pptx"
            self._write_presentation(pptx_path)

            errors = "\n".join(
                verify_presentation(pptx_path, source_path)
            ).lower()

        self.assertIn("source", errors)
        self.assertIn("14", errors)

    def test_rejects_out_of_sequence_source_numbers(self) -> None:
        with TemporaryDirectory() as directory:
            source_path = Path(directory) / "source.md"
            source_path.write_text(
                SOURCE_PATH.read_text(encoding="utf-8").replace(
                    "## 第 14 頁", "## 第 13 頁", 1
                ),
                encoding="utf-8",
            )
            pptx_path = Path(directory) / "deck.pptx"
            self._write_presentation(pptx_path)

            errors = "\n".join(
                verify_presentation(pptx_path, source_path)
            ).lower()

        self.assertIn("source", errors)
        self.assertIn("1..14", errors)

    def test_rejects_non_widescreen_presentation(self) -> None:
        self.assertIn("16:9", self._errors(aspect_ratio="standard"))

    def test_rejects_slide_without_named_conclusion(self) -> None:
        self.assertIn("conclusion", self._errors(missing_conclusion_slide=3))

    def test_rejects_slide_with_fewer_than_three_named_bullets(self) -> None:
        self.assertIn("3–5", self._errors(bullet_count_by_slide={4: 2}))

    def test_rejects_slide_with_more_than_five_named_bullets(self) -> None:
        self.assertIn("3–5", self._errors(bullet_count_by_slide={4: 6}))

    def test_rejects_wrong_deep_blue_background(self) -> None:
        self.assertIn("081a2e", self._errors(background_color="FFFFFF"))

    def test_rejects_obviously_nonwhite_primary_text(self) -> None:
        self.assertIn("white", self._errors(primary_text_color="000000"))

    def test_rejects_deck_without_teal_accent(self) -> None:
        self.assertIn("21d4b4", self._errors(accent_color="FF0000"))

    def test_rejects_second_slide_without_teal_accent(self) -> None:
        self.assertIn("slide 2", self._errors(missing_accent_slide=2))

    def test_rejects_second_slide_with_off_canvas_teal_accent(self) -> None:
        self.assertIn(
            "21d4b4", self._errors(invisible_accent_slide=2)
        )

    def test_rejects_wrong_primary_font(self) -> None:
        self.assertIn(
            "microsoft jhenghei", self._errors(primary_font_name="Arial")
        )

    def test_rejects_unnamed_visible_black_arial_text(self) -> None:
        errors = self._errors(include_unnamed_wrong_style=True)
        self.assertIn("white", errors)
        self.assertIn("microsoft jhenghei", errors)

    def test_rejects_one_slide_with_wrong_primary_text_theme(self) -> None:
        errors = self._errors(wrong_style_slide=2)
        self.assertIn("slide 2", errors)
        self.assertIn("white", errors)
        self.assertIn("microsoft jhenghei", errors)

    def test_rejects_english_bullet_substitution(self) -> None:
        self.assertIn(
            "source bullet",
            self._errors(bullet_text_override="English replacement bullet"),
        )

    def test_rejects_english_conclusion_substitution(self) -> None:
        self.assertIn(
            "source conclusion",
            self._errors(conclusion_text_override="English conclusion"),
        )

    def test_rejects_english_notes_substitution(self) -> None:
        self.assertIn(
            "source speaker notes",
            self._errors(notes_text_override="English speaker guidance"),
        )

    def test_rejects_missing_flow_node_shape(self) -> None:
        self.assertIn("flow-node", self._errors(flow_node_count=5))

    def test_rejects_flow_node_picture_replacement(self) -> None:
        self.assertIn(
            "native", self._errors(replace_flow_node_with_picture=True)
        )

    def test_rejects_english_flow_node_substitution(self) -> None:
        self.assertIn(
            "source flow node",
            self._errors(flow_node_text_override="English Node"),
        )

    def test_rejects_missing_technical_card(self) -> None:
        self.assertIn("exactly 13", self._errors(technology_count=12))

    def test_rejects_technical_card_that_does_not_match_source(self) -> None:
        self.assertIn(
            "technical name",
            self._errors(technology_name_override="Unknown Technology"),
        )

    def test_rejects_technical_card_picture_replacement(self) -> None:
        self.assertIn(
            "native", self._errors(replace_tech_card_with_picture=True)
        )

    def test_rejects_technical_card_with_tampered_status(self) -> None:
        self.assertIn(
            "status",
            self._errors(technology_status_override="Status tampered"),
        )

    def test_rejects_technical_row_with_tampered_value(self) -> None:
        self.assertIn(
            "value",
            self._errors(technology_value_override="Value tampered"),
        )

    def test_rejects_undersized_conclusion_text(self) -> None:
        self.assertIn(
            "16pt",
            self._errors(undersized_shape_name="conclusion-2"),
        )

    def test_rejects_undersized_bullet_text(self) -> None:
        self.assertIn(
            "16pt",
            self._errors(undersized_shape_name="bullet-2-1"),
        )

    def test_rejects_undersized_technical_card_text(self) -> None:
        source_technology = parse_markdown(SOURCE_PATH)[8].table_rows[0].cells[0]
        self.assertIn(
            "12pt",
            self._errors(
                undersized_shape_name=(
                    f"tech-card-{self._normalized_name(source_technology)}"
                )
            ),
        )

    def test_rejects_undersized_roadmap_step_text(self) -> None:
        self.assertIn(
            "16pt",
            self._errors(undersized_shape_name="roadmap-step-1"),
        )

    def test_rejects_tampered_visible_status_label(self) -> None:
        self.assertIn(
            "status-tag",
            self._errors(status_label_override="Tampered status"),
        )

    def test_rejects_missing_named_roadmap_step(self) -> None:
        names = [f"roadmap-step-{number}" for number in range(1, 5)]
        self.assertIn("roadmap-step-5", self._errors(roadmap_names=names))

    def test_rejects_duplicate_named_roadmap_step(self) -> None:
        names = [f"roadmap-step-{number}" for number in range(1, 6)]
        names.append("roadmap-step-5")
        self.assertIn("exactly once", self._errors(roadmap_names=names))

    def test_rejects_roadmap_step_picture_replacement(self) -> None:
        self.assertIn(
            "native", self._errors(replace_roadmap_step_with_picture=True)
        )

    def test_rejects_english_roadmap_step_substitution(self) -> None:
        self.assertIn(
            "source roadmap",
            self._errors(roadmap_text_override="1. English Step"),
        )

    def test_rejects_sensitive_text_inside_nested_group(self) -> None:
        self.assertIn(
            "sensitive",
            self._errors(include_grouped_sensitive_text=True),
        )

    def test_rejects_sensitive_table_cell_inside_group(self) -> None:
        self.assertIn(
            "sensitive",
            self._errors(include_grouped_sensitive_table=True),
        )

    def test_rejects_title_mismatch(self) -> None:
        self.assertIn("title", self._errors(title_override="Unexpected title"))

    def test_rejects_empty_speaker_notes(self) -> None:
        self.assertIn("speaker notes", self._errors(include_notes=False))

    def test_rejects_out_of_bounds_shape(self) -> None:
        self.assertIn(
            "outside", self._errors(include_out_of_bounds_shape=True)
        )

    def test_rejects_top_level_sensitive_identifier(self) -> None:
        self.assertIn(
            "sensitive", self._errors(include_sensitive_text=True)
        )


if __name__ == "__main__":
    unittest.main()
