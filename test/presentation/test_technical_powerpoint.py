"""Contract tests for the technical presentation verifier."""

from __future__ import annotations

import hashlib
from io import BytesIO
from pathlib import Path
from tempfile import TemporaryDirectory
import unittest
import zipfile

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from scripts.presentation.build_client_powerpoint import parse_markdown
from scripts.presentation.build_technical_powerpoint import build_technical_presentation
from scripts.presentation.verify_technical_powerpoint import verify_technical_presentation


ROOT = Path(__file__).resolve().parents[2]
SOURCE_PATH = ROOT / "docs/presentations/2026-07-30-technical-achievements-presentation.md"
ONE_PIXEL_PNG = b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x06\x00\x00\x00\x1f\x15\xc4\x89\x00\x00\x00\x0dIDAT\x08\xd7c\xf8\xcf\xc0\xf0\x1f\x00\x05\x00\x01\xff\x89\x99=\x1d\x00\x00\x00\x00IEND\xaeB`\x82"


def _roadmap_label(bullet: str) -> str:
    number, body = bullet.split(". ", 1)
    return f"{number}. {body.split('：', 1)[0]}"


def _display_text(text: str) -> str:
    return text.replace("**", "").replace("__", "").replace("`", "")


def _triangle_tail(connector) -> None:
    connector.line.width = Pt(1)
    tail = OxmlElement("a:tailEnd")
    tail.set("type", "triangle")
    connector._element.spPr.ln.append(tail)


class TechnicalSourceTests(unittest.TestCase):
    def test_source_has_sixteen_sequential_traditional_chinese_slides(self) -> None:
        slides = parse_markdown(SOURCE_PATH)
        self.assertEqual(list(range(1, 17)), [slide.number for slide in slides])
        self.assertTrue(all(slide.title.strip() for slide in slides))
        self.assertTrue(all(3 <= len(slide.bullets) <= 5 for slide in slides))
        self.assertTrue(all(slide.speaker_notes.strip() for slide in slides))
        self.assertTrue(all("開發中" in " ".join(slide.bullets) for slide in slides[12:13]))
        self.assertIn("待合併前驗證", " ".join(slides[13].bullets))


class TechnicalVerifierContractTests(unittest.TestCase):
    def _add_text(self, shapes, name: str, text: str) -> None:
        shape = shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(10), Inches(0.65))
        shape.name = name
        shape.text = text
        if text:
            run = shape.text_frame.paragraphs[0].runs[0]
            run.font.name = "Microsoft JhengHei"
            run.font.size = Pt(16)
            run.font.color.rgb = RGBColor.from_string("F4F8FC")

    def _write_deck(
        self, destination: Path, *, title_override: str | None = None,
        bullet_override: str | None = None, conclusion_override: str | None = None,
        notes_override: str | None = None, duplicate_name: str | None = None,
        empty_special: tuple[int, str] | None = None, omit_special: tuple[int, str] | None = None,
        sensitive_location: str | None = None, quality_override: str | None = None,
        maturity_override: str | None = None, roadmap_override: str | None = None,
        rename_special: tuple[int, str, str] | None = None,
        message_connector_replacement: str | None = None,
        duplicate_message_connector: bool = False,
        narrow_bullet: bool = False,
        architecture_role_override: str | None = None,
        undersized_shape: str | None = None,
        missing_knowledge_connector: int | None = None,
    ) -> None:
        presentation = Presentation()
        presentation.slides._sldIdLst.clear()
        presentation.slide_width = Inches(40 / 3)
        presentation.slide_height = Inches(7.5)
        source_slides = parse_markdown(SOURCE_PATH)
        for source in source_slides:
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            background = slide.background.fill
            background.solid()
            background.fore_color.rgb = RGBColor.from_string("081A2E")
            accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(40 / 3), Inches(.08))
            accent.name = f"accent-{source.number}"
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor.from_string("21D4B4")
            accent.line.fill.background()
            self._add_text(slide.shapes, f"title-{source.number}", title_override if source.number == 1 and title_override is not None else source.title)
            self._add_text(slide.shapes, f"conclusion-{source.number}", conclusion_override if source.number == 1 and conclusion_override is not None else _display_text(source.bullets[0]))
            for index, bullet in enumerate(source.bullets, start=1):
                text = bullet_override if source.number == 1 and index == 1 and bullet_override is not None else _display_text(bullet)
                self._add_text(slide.shapes, f"bullet-{source.number}-{index}", text)
            self._add_text(slide.shapes, f"footer-{source.number}", "footer")
            self._add_text(slide.shapes, f"page-number-{source.number}", str(source.number))
            slide.notes_slide.notes_text_frame.text = notes_override if source.number == 1 and notes_override is not None else source.speaker_notes

        def special(slide_number: int, name: str, text: str) -> None:
            if omit_special == (slide_number, name):
                return
            output_name = rename_special[2] if rename_special and rename_special[:2] == (slide_number, name) else name
            self._add_text(presentation.slides[slide_number - 1].shapes, output_name, "" if empty_special == (slide_number, name) else text)

        for index, role in enumerate(("使用者", "LINE", "Worker", "Queue", "AI", "D1", "Open-Meteo"), start=1):
            special(3, f"architecture-node-{index}", architecture_role_override if index == 1 and architecture_role_override is not None else role)
        for index in range(6):
            connector = presentation.slides[2].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
            connector.name = f"architecture-connector-{index + 1}"
            _triangle_tail(connector)
        for index, text in enumerate(("Webhook 接收", "簽章與群組 mention 驗證", "Queue 入列", "AI 或資料來源處理", "LINE reply 與 push fallback", "D1 與觀測紀錄"), start=1):
            special(4, f"message-flow-step-{index}", text)
        for index in range(5):
            name = f"message-flow-connector-{index + 1}"
            if index == 0 and message_connector_replacement == "textbox":
                self._add_text(presentation.slides[3].shapes, name, "假 connector")
            elif index == 0 and message_connector_replacement == "picture":
                shape = presentation.slides[3].shapes.add_picture(BytesIO(ONE_PIXEL_PNG), Inches(1), Inches(1), Inches(1), Inches(.3))
                shape.name = name
            else:
                connector = presentation.slides[3].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
                connector.name = name
                _triangle_tail(connector)
        if duplicate_message_connector:
            connector = presentation.slides[3].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
            connector.name = "message-flow-connector-1-copy"
            _triangle_tail(connector)
        for index in range(3):
            connector = presentation.slides[4].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
            connector.name = f"worker-flow-connector-{index + 1}"
            _triangle_tail(connector)
        for name in ("retry", "dlq", "deduplication"):
            special(6, f"reliability-node-{name}", name)
        for index in range(2):
            connector = presentation.slides[5].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
            connector.name = f"reliability-connector-{index + 1}"
            _triangle_tail(connector)
        for name, text in (("question-record", "D1 問題紀錄"), ("weather-cache", "weather cache"), ("group-settings", "group settings"), ("metrics", "metrics"), ("lifecycle", "lifecycle")):
            special(9, f"data-node-{name}", text)
        special(11, "observability-correlation-webhook", "webhookEventId")
        special(11, "observability-correlation-operation", "operationId")
        for index, text in enumerate(("webhook.enqueue.completed", "question.started", "storage.claim.completed", "answer.completed", "line.reply.completed"), start=1):
            special(11, f"observability-event-{index}", text)
        for index in range(1, 5):
            connector = presentation.slides[10].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
            connector.name = f"observability-connector-{index}"
            _triangle_tail(connector)
        for name, text in (("main", "主線 134 通過"), ("knowledge", "知識搜尋 421 通過"), ("timeout", "1 項超過 5 秒"), ("predeploy", "設定檢查待更新")):
            special(12, f"quality-gate-{name}", quality_override if name == "main" and quality_override is not None else text)
        for name in ("r2", "ingestion-queue", "workers-ai", "vectorize", "retrieval", "grounded-answer"):
            special(13, f"knowledge-node-{name}", name)
        for index in range(1, 6):
            if missing_knowledge_connector == index:
                continue
            connector = presentation.slides[12].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(index), Inches(1), Inches(index + 1), Inches(1))
            connector.name = f"knowledge-connector-{index}"
            _triangle_tail(connector)
        special(13, "development-status-13", "開發中")
        special(14, "maturity-matrix-14", maturity_override or "已完成｜開發中｜待合併前驗證｜待正式環境驗證")
        for index in range(1, 6):
            special(16, f"roadmap-step-{index}", roadmap_override if index == 1 and roadmap_override is not None else _roadmap_label(source_slides[15].bullets[index - 1]))
        for index in range(1, 5):
            connector = presentation.slides[15].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
            connector.name = f"roadmap-connector-{index}"
            _triangle_tail(connector)
        for name, text in (("primary", "主要模型"), ("fallback", "備援模型"), ("policy", "回答政策")):
            special(7, f"model-boundary-{name}", text)
        for name, text in (("intent", "意圖辨識"), ("cache-read", "快取讀取"), ("provider", "Open-Meteo"), ("cache-write", "快取寫入")):
            special(8, f"weather-cache-step-{name}", text)
        for index in range(1, 4):
            connector = presentation.slides[7].shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(2), Inches(1))
            connector.name = f"weather-cache-connector-{index}"
            _triangle_tail(connector)
        for name, text in (("logs", "Workers Logs"), ("d1", "D1"), ("secrets", "Cloudflare secrets")):
            special(10, f"privacy-layer-{name}", text)
        if duplicate_name:
            self._add_text(presentation.slides[0].shapes, duplicate_name, "重複名稱")
        if narrow_bullet:
            narrow = next(
                shape
                for shape in presentation.slides[0].shapes
                if shape.name == "bullet-1-1"
            )
            narrow.width = Inches(0.25)
        if undersized_shape:
            shape = next(
                shape
                for slide in presentation.slides
                for shape in slide.shapes
                if shape.name == undersized_shape
            )
            shape.text_frame.paragraphs[0].runs[0].font.size = Pt(15)
        if sensitive_location == "top":
            self._add_text(presentation.slides[0].shapes, "sensitive-top", "access_token")
        elif sensitive_location == "group":
            group = presentation.slides[0].shapes.add_group_shape()
            inner = group.shapes.add_group_shape()
            self._add_text(inner.shapes, "sensitive-group", "channel_secret")
        elif sensitive_location == "table":
            group = presentation.slides[0].shapes.add_group_shape()
            table = presentation.slides[0].shapes.add_table(1, 1, Inches(1), Inches(1), Inches(2), Inches(.4))
            table.table.cell(0, 0).text = "database_id"
            group.shapes._spTree.insert_element_before(table._element, "p:extLst")
        elif sensitive_location == "notes":
            presentation.slides[0].notes_slide.notes_text_frame.text = "analytics_hash_key"
        presentation.save(destination)

    def _errors(self, **kwargs) -> str:
        with TemporaryDirectory() as directory:
            path = Path(directory) / "deck.pptx"
            self._write_deck(path, **kwargs)
            return "\n".join(verify_technical_presentation(path, SOURCE_PATH)).lower()

    def test_accepts_exact_source_bound_deck(self) -> None:
        self.assertEqual("", self._errors())

    def test_rejects_title_not_equal_to_source(self) -> None:
        self.assertIn("title", self._errors(title_override="English title"))

    def test_rejects_bullet_not_equal_to_source(self) -> None:
        self.assertIn("bullet", self._errors(bullet_override="English bullet"))

    def test_rejects_conclusion_not_equal_to_first_source_bullet(self) -> None:
        self.assertIn("conclusion", self._errors(conclusion_override="English conclusion"))

    def test_rejects_notes_without_source_speaker_notes(self) -> None:
        self.assertIn("speaker notes", self._errors(notes_override="English notes"))

    def test_rejects_copy_suffix_shape_name(self) -> None:
        self.assertIn("title-1", self._errors(duplicate_name="title-1-copy"))

    def test_rejects_sensitive_top_level_text(self) -> None:
        self.assertIn("sensitive", self._errors(sensitive_location="top"))

    def test_rejects_sensitive_nested_group_text(self) -> None:
        self.assertIn("sensitive", self._errors(sensitive_location="group"))

    def test_rejects_sensitive_table_cell(self) -> None:
        self.assertIn("sensitive", self._errors(sensitive_location="table"))

    def test_rejects_sensitive_speaker_notes(self) -> None:
        self.assertIn("sensitive", self._errors(sensitive_location="notes"))

    def test_rejects_empty_architecture_node(self) -> None:
        self.assertIn("architecture-node", self._errors(empty_special=(3, "architecture-node-1")))

    def test_rejects_renamed_reliability_node(self) -> None:
        self.assertIn("reliability-node-retry", self._errors(rename_special=(6, "reliability-node-retry", "reliability-node-retry-copy")))

    def test_rejects_missing_message_flow_node(self) -> None:
        self.assertIn("message-flow-step", self._errors(omit_special=(4, "message-flow-step-1")))

    def test_rejects_empty_reliability_node(self) -> None:
        self.assertIn("reliability-node", self._errors(empty_special=(6, "reliability-node-retry")))

    def test_rejects_empty_data_node(self) -> None:
        self.assertIn("data-node", self._errors(empty_special=(9, "data-node-metrics")))

    def test_rejects_renamed_data_node(self) -> None:
        self.assertIn("data-node-question-record", self._errors(rename_special=(9, "data-node-question-record", "data-node-question-record-copy")))

    def test_rejects_empty_observability_event(self) -> None:
        self.assertIn("observability-event", self._errors(empty_special=(11, "observability-event-1")))

    def test_rejects_quality_failure_claim(self) -> None:
        self.assertIn("quality-gate", self._errors(quality_override="主線 134 失敗"))

    def test_rejects_missing_development_status_text(self) -> None:
        self.assertIn("development-status", self._errors(empty_special=(13, "development-status-13")))

    def test_rejects_renamed_knowledge_node(self) -> None:
        self.assertIn("knowledge-node-r2", self._errors(rename_special=(13, "knowledge-node-r2", "knowledge-node-r2-copy")))

    def test_rejects_maturity_matrix_missing_required_state(self) -> None:
        self.assertIn("maturity", self._errors(maturity_override="已完成｜開發中｜待驗證"))

    def test_rejects_maturity_matrix_copy_suffix(self) -> None:
        self.assertIn("maturity-matrix-14", self._errors(rename_special=(14, "maturity-matrix-14", "maturity-matrix-14-copy")))

    def test_rejects_percentage_maturity_matrix(self) -> None:
        self.assertIn("percentage", self._errors(maturity_override="已完成 80%｜開發中"))

    def test_rejects_english_roadmap_node(self) -> None:
        self.assertIn("roadmap-step-1", self._errors(roadmap_override="1. English step"))

    def test_rejects_textbox_instead_of_message_flow_connector(self) -> None:
        self.assertIn("message-flow-connector-1", self._errors(message_connector_replacement="textbox"))

    def test_rejects_non_line_message_flow_connector(self) -> None:
        self.assertIn("message-flow-connector-1", self._errors(message_connector_replacement="picture"))

    def test_rejects_message_flow_connector_copy_suffix(self) -> None:
        self.assertIn("message-flow-connector-1", self._errors(duplicate_message_connector=True))

    def test_rejects_a_narrow_long_source_bullet_card(self) -> None:
        self.assertIn("density", self._errors(narrow_bullet=True))

    def test_rejects_an_architecture_role_replaced_with_webhook(self) -> None:
        self.assertIn("architecture-node", self._errors(architecture_role_override="Webhook"))

    def test_rejects_an_architecture_node_below_sixteen_points(self) -> None:
        self.assertIn("16pt", self._errors(undersized_shape="architecture-node-1"))

    def test_rejects_a_missing_knowledge_chain_connector(self) -> None:
        self.assertIn("knowledge-connector-3", self._errors(missing_knowledge_connector=3))

    def test_rejects_picture_anywhere_in_the_deck(self) -> None:
        self.assertIn(
            "image",
            self._errors(message_connector_replacement="picture"),
        )

    def test_rejects_external_ooxml_relationship(self) -> None:
        with TemporaryDirectory() as directory:
            path = Path(directory) / "deck.pptx"
            self._write_deck(path)
            with zipfile.ZipFile(path, "a") as archive:
                archive.writestr(
                    "ppt/slides/_rels/slide1.xml.rels",
                    '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"><Relationship Id="rId99" Type="urn:test" Target="https://example.test" TargetMode="External"/></Relationships>',
                )
            errors = "\n".join(verify_technical_presentation(path, SOURCE_PATH)).lower()
        self.assertIn("external", errors)

    def test_rejects_vba_or_media_parts(self) -> None:
        with TemporaryDirectory() as directory:
            path = Path(directory) / "deck.pptx"
            self._write_deck(path)
            with zipfile.ZipFile(path, "a") as archive:
                archive.writestr("ppt/vbaProject.bin", b"not a macro")
                archive.writestr("ppt/media/clip.mp4", b"not media")
            errors = "\n".join(verify_technical_presentation(path, SOURCE_PATH)).lower()
        self.assertIn("vba", errors)
        self.assertIn("media", errors)

    def test_rejects_invalid_relationship_xml(self) -> None:
        with TemporaryDirectory() as directory:
            path = Path(directory) / "deck.pptx"
            self._write_deck(path)
            with zipfile.ZipFile(path, "a") as archive:
                archive.writestr("ppt/slides/_rels/slide1.xml.rels", "<Relationships>")
            errors = "\n".join(verify_technical_presentation(path, SOURCE_PATH)).lower()
        self.assertIn("ooxml", errors)


class GeneratedTechnicalPowerPointTests(unittest.TestCase):
    def test_generated_deck_satisfies_the_complete_technical_contract(self) -> None:
        with TemporaryDirectory() as directory:
            output_path = Path(directory) / "technical.pptx"

            build_technical_presentation(SOURCE_PATH, output_path)

            self.assertEqual(
                [], verify_technical_presentation(output_path, SOURCE_PATH)
            )

    def test_generated_deck_has_safe_theme_bounds_and_directed_connectors(self) -> None:
        client_path = ROOT / "docs/presentations/AI-ChotBot-project-progress-client.pptx"
        before_hash = hashlib.sha256(client_path.read_bytes()).hexdigest()
        before_mtime = client_path.stat().st_mtime_ns
        with TemporaryDirectory() as directory:
            output_path = Path(directory) / "technical.pptx"
            build_technical_presentation(SOURCE_PATH, output_path)
            deck = Presentation(output_path)

        self.assertAlmostEqual(13.333, deck.slide_width.inches, places=3)
        self.assertEqual(7.5, deck.slide_height.inches)
        for index, slide in enumerate(deck.slides, start=1):
            self.assertEqual("081A2E", str(slide.background.fill.fore_color.rgb))
            accent = next(shape for shape in slide.shapes if shape.name == f"accent-{index}")
            self.assertEqual("21D4B4", str(accent.fill.fore_color.rgb))
            for shape in slide.shapes:
                if shape.name.startswith(("accent-", "footer-", "page-number-", "section-")):
                    continue
                self.assertLessEqual(shape.left + shape.width, deck.slide_width)
                self.assertLessEqual(shape.top + shape.height, Inches(7.0))
        for slide_number, prefix in ((3, "architecture-connector-"), (4, "message-flow-connector-"), (5, "worker-flow-connector-"), (6, "reliability-connector-"), (11, "observability-connector-"), (16, "roadmap-connector-")):
            for shape in deck.slides[slide_number - 1].shapes:
                if shape.name.startswith(prefix):
                    tail = shape._element.spPr.ln.find(qn("a:tailEnd"))
                    self.assertIsNotNone(tail, f"{shape.name} requires a tail arrow")
                    self.assertEqual("triangle", tail.get("type"))
        self.assertEqual(before_hash, hashlib.sha256(client_path.read_bytes()).hexdigest())
        self.assertEqual(before_mtime, client_path.stat().st_mtime_ns)

    def test_generated_deck_has_semantic_special_visuals_and_short_roadmap_labels(self) -> None:
        with TemporaryDirectory() as directory:
            output_path = Path(directory) / "technical.pptx"
            build_technical_presentation(SOURCE_PATH, output_path)
            deck = Presentation(output_path)

        flow = deck.slides[3]
        expected_flow = ("Webhook 接收", "簽章與群組 mention 驗證", "Queue 入列", "AI 或資料來源處理", "LINE reply 與 push fallback", "D1 與觀測紀錄")
        self.assertEqual(expected_flow, tuple(next(shape.text for shape in flow.shapes if shape.name == f"message-flow-step-{index}") for index in range(1, 7)))
        events = ("webhook.enqueue.completed", "question.started", "storage.claim.completed", "answer.completed", "line.reply.completed")
        self.assertEqual(events, tuple(next(shape.text for shape in deck.slides[10].shapes if shape.name == f"observability-event-{index}") for index in range(1, 6)))
        for slide_number, names in ((7, ("model-boundary-primary", "model-boundary-fallback", "model-boundary-policy")), (8, ("weather-cache-step-intent", "weather-cache-step-cache-read", "weather-cache-step-provider", "weather-cache-step-cache-write")), (10, ("privacy-layer-logs", "privacy-layer-d1", "privacy-layer-secrets"))):
            actual_names = {shape.name for shape in deck.slides[slide_number - 1].shapes}
            self.assertTrue(set(names).issubset(actual_names))
        source = parse_markdown(SOURCE_PATH)
        self.assertEqual(tuple(_roadmap_label(bullet) for bullet in source[15].bullets), tuple(next(shape.text for shape in deck.slides[15].shapes if shape.name == f"roadmap-step-{index}") for index in range(1, 6)))

    def test_generated_special_bullet_cards_are_readable_and_strip_inline_markdown(self) -> None:
        with TemporaryDirectory() as directory:
            output_path = Path(directory) / "technical.pptx"
            build_technical_presentation(SOURCE_PATH, output_path)
            deck = Presentation(output_path)

        for slide_number in (3, 4, 5, 6, 7, 8, 9, 10, 11, 12, 13, 14, 16):
            bullets = [
                shape
                for shape in deck.slides[slide_number - 1].shapes
                if shape.name.startswith(f"bullet-{slide_number}-")
            ]
            self.assertLessEqual(len({shape.left for shape in bullets}), 2)
            self.assertTrue(all(shape.width.inches >= 5.8 for shape in bullets))
            self.assertTrue(all(shape.height.inches >= 0.65 for shape in bullets))
            self.assertTrue(all(shape.top.inches + shape.height.inches <= 7.0 for shape in bullets))
            self.assertTrue(all(len(shape.text) / (shape.width.inches * shape.height.inches) <= 25 for shape in bullets))
        development_bullet = next(
            shape.text
            for shape in deck.slides[12].shapes
            if shape.name == "bullet-13-1"
        )
        self.assertNotIn("**", development_bullet)

    def test_generated_architecture_roles_and_knowledge_chain_are_exact(self) -> None:
        with TemporaryDirectory() as directory:
            output_path = Path(directory) / "technical.pptx"
            build_technical_presentation(SOURCE_PATH, output_path)
            deck = Presentation(output_path)

        architecture = deck.slides[2]
        expected_roles = ("使用者", "LINE", "Worker", "Queue", "AI", "D1", "Open-Meteo")
        self.assertEqual(
            expected_roles,
            tuple(
                next(
                    shape.text
                    for shape in architecture.shapes
                    if shape.name == f"architecture-node-{index}"
                )
                for index in range(1, 8)
            ),
        )
        for index in range(1, 8):
            shape = next(shape for shape in architecture.shapes if shape.name == f"architecture-node-{index}")
            self.assertGreaterEqual(shape.text_frame.paragraphs[0].runs[0].font.size.pt, 16)
        knowledge = deck.slides[12]
        for index in range(1, 6):
            connector = next(shape for shape in knowledge.shapes if shape.name == f"knowledge-connector-{index}")
            self.assertEqual(MSO_SHAPE_TYPE.LINE, connector.shape_type)
            self.assertLess(connector.begin_x, connector.end_x)
            tail = connector._element.spPr.ln.find(qn("a:tailEnd"))
            self.assertIsNotNone(tail)
            self.assertEqual("triangle", tail.get("type"))


if __name__ == "__main__":
    unittest.main()
